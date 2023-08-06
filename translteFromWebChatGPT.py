import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pygoogletranslation import Translator
from deep_translator import GoogleTranslator
from pptx.util import Inches
import mysql.connector
import tkinter as tk
from tkinter import *
import random
import os
import openai
import time

# Function to translate text using ChatGPT
openai.api_key = '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'
def translate_text(text, source_language, target_language):
    prompt = f"Translate the following '{source_language}' text to '{target_language}': {text}"

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that translates text."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=150,
        n=1,
        stop=None,
        temperature=0.5,
    )

    translation = response.choices[0].message.content.strip()
    return translation


# GUI Creation
def button_click():
    global input_topic, input_count
    input_topic = topic_field.get()
    input_count = int(count_field.get())
    label.config(text="Input retrieved.")

    #Check translation selection
    if (variable.get() == "Afrikaans"):
         print("Translate language is "+variable.get())
         variable.set("af")
         languageOutput.set("Afrikaans")

    elif (variable.get() == "Telugu"):
         print("Translate language is "+variable.get())
         variable.set("te")
         languageOutput.set("Telugu")

    elif (variable.get() == "Hindi"):
         print("Translate language is "+variable.get())
         variable.set("hi")
         languageOutput.set("Hindi")

    elif (variable.get() == "Chinese"):             #Chinese (simplified)': 'zh-CN', Chinese (traditional)': 'zh-TW'
         print("Translate language is "+variable.get())
         variable.set("zh-CN")
         languageOutput.set("Chinese")

    elif (variable.get() == "Urdu"):
         print("Translate language is "+variable.get())
         variable.set("ur")
         languageOutput.set("Urdu")
         
    elif (variable.get() == "German"):
         print("Translate language is "+variable.get())
         variable.set("de")
         languageOutput.set("German")
    
    window.destroy()
    

window = tk.Tk()
window.resizable(height=400, width=400)
window.geometry('800x800')
window['background'] = 'white'
window.title('QuizMaster using ChatGPT')

label = tk.Label(window, text="Welcome to QuizMaster", font = ("Arial", 40), fg = "black", bg = "white")
topic_label = tk.Label(window, text="Please enter a topic:", font = ("Arial", 25), fg = "black", bg = "white")
count_label = tk.Label(window, text="How many questions should be chosen?", font = ("Arial", 25), fg = "black", bg = "white")
topic_field = tk.Entry(window, width = 25)
count_field = tk.Entry(window, width = 25)

#Translating options
OPTIONS = [
"Afrikaans",
"Telugu",
"Hindi",
"Chinese",
"Urdu",
"German"
] 

variable = StringVar(window)
languageOutput = StringVar(window)
variable.set(OPTIONS[0]) # default value
w = OptionMenu(window, variable, *OPTIONS)
w.pack()


button = tk.Button(window, text="Submit",background= "blue", command=button_click ,font = ("Arial", 25), width = 8)

label.pack()
topic_label.pack()
topic_field.pack()
count_label.pack()
count_field.pack()
button.pack()

window.mainloop()

# Connecting to database
conn = mysql.connector.connect(
    host='localhost',
    user='root',
    #password='password',
    database='quiz_master'
)

#Query that's being executed
cursor = conn.cursor()
query = "SELECT * FROM questions WHERE topic = '" + input_topic + "'"
cursor.execute(query)

results = cursor.fetchall()
if (len(results) < input_count):
      input_count = len(results)

random_results = random.sample(results, input_count)

# Create a Presentation object
presentation = Presentation()
slide_layout = presentation.slide_layouts[1]
slide_count = 0

#Go through questions and questin titles
for row in random_results:
        slide = presentation.slides.add_slide(slide_layout)
        slideTranslate = presentation.slides.add_slide(slide_layout)
        translateShapes = slideTranslate.shapes
        shapes = slide.shapes

# Set the question as the slide title
        title_shape = shapes.title
        titleTrans_shape = translateShapes.title
        question = row[2]
        questionTranslate = translate_text(question, 'auto', variable.get())
        # time.sleep(20)
        title_shape.text = question
        titleTrans_shape.text = questionTranslate


        #TRANSLATE USING CHATGPT
        answers = [row[3], row[4], row[5], row[6]]
        bodyTrans_shape = translateShapes.placeholders[1]
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tq = bodyTrans_shape.text_frame
        t0 = translate_text(answers[0], 'auto', variable.get())
        tf.text = "A. " + answers[0]
        tq.text = "A. " + t0

        # p = tf.add_paragraph()
        # t1=translate_text(answers[1], 'auto', variable.get())
        # tf.text = "A. " + answers[1]
        # tq.text = "A. " + t1
        

        # p = tf.add_paragraph()
        # t2=translate_text(answers[2], 'auto', variable.get())
        # tf.text = "A. " + answers[2] + " <--"+languageOutput.get() +" Translation --> " + ''.join([str(t2)])
        # time.sleep(20)

        # p = tf.add_paragraph()
        # t3=translate_text(answers[3], 'auto', variable.get())
        # tf.text = "A. " + answers[3] + " <--"+languageOutput.get() +" Translation --> " + ''.join([str(t3)])



        #TRANSLATE USING GOOGLE TRANSLATE

        #Last 3 answers are translated using google translate instead of ChatGPT because  
        #rate limit is reached for default-gpt-3.5-turbo 

        # p = tf.add_paragraph()
        # t0=GoogleTranslator(source='auto', target= variable.get()).translate(answers[0])
        # p.text = "C. " + answers[0] +" <--"+languageOutput.get() + " Translation --> " + ''.join([str(t0)])

        t = tq.add_paragraph()
        p = tf.add_paragraph()
        t1 = GoogleTranslator(source='auto', target= variable.get()).translate(answers[1])
        p.text = "B. " + answers[1] 
        t.text = "B. " + t1

        t = tq.add_paragraph()
        p = tf.add_paragraph()
        t2=GoogleTranslator(source='auto', target= variable.get()).translate(answers[2])
        p.text = "C. " + answers[2] 
        t.text = "C. " + t2

        t = tq.add_paragraph()
        p = tf.add_paragraph()
        t3 = GoogleTranslator(source='auto', target= variable.get()).translate(answers[3])
        p.text = "D. " + answers[3] 
        t.text = "D. " + t3

        slide_count += 1

        print(t0)
        print(t1)
        print(t2)
        print(t3)

# Save the PowerPoint presentation to a file
presentation.save(r"C:\Users\coenr\Downloads\QuizmasterPowerpoint-main\QuizmasterPowerpoint-main\output.pptx")


cursor.close()
conn.close()