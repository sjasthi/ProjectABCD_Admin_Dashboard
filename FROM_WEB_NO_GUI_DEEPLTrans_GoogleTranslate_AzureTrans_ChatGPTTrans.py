import os
from io import open
import shutil
import collections
import collections.abc
from os.path import basename
import pptx.util
from pptx import Presentation
from pptx.util import Pt
from bs4 import BeautifulSoup
import requests
import deepl
from pygoogletranslation import Translator
from deep_translator import GoogleTranslator
import openai
import time

""" 
This program uses different APIs to do a text translation and creates a powerpoint slides to show the output.

APIs: Google Translate , Azure Translate, ChatGPT Translate, and DeepL Translate.

Text translation using DeepL API!
-To use the DeepL Python Library, you'll need an API authentication key. 
To get a key, please create an account here -->https://www.deepl.com/docs-api/api-access/
With a DeepL API Free account you can translate up to 500,000 characters/month for free.

Installation
The library can be installed from PyPI using pip:

pip install --upgrade deepl """

methods = "Web"
slideOption = 1
textFont = titleFont = "NATS"
textSize = 13
titleSize = 2
output_file = 'project_abcd.pptx'
dest_language = "ar"  # Arabic language
dest_language_DeepL="FR" # French language

TranslationEngines = {}
TranslationEngines["DeepLAPI"] = "DeepL"
TranslationEngines['AzureAPI'] = 'Azure'
TranslationEngines['GoogleAPI'] = 'Google'
TranslationEngines['ChatGPTAPI'] = 'ChatGPT'

#function to provide users with a list of options to choose from
def selectFromDict(options, Engine):
    index = 0

    indexValidList = []
    print('Select an ' + Engine + ':')
    for optionName in options:
        index = index + 1
        indexValidList.extend([options[optionName]])
        print(str(index) + ') ' + optionName)
        inputValid = False
    while not inputValid:
        inputRaw = input(Engine + ': ')
        inputNo = int(inputRaw) - 1
        if inputNo > -1 and inputNo < len(indexValidList):
            selected = indexValidList[inputNo]
            print('Selected ' + Engine + ': ' + selected)
            inputValid = True
            break
        else:
            print('Please select a valid ' + Engine + ' number')


    return selected

    print(selected)

inp = str(selectFromDict(TranslationEngines, 'Engine'))

#print("input: ", inp)

""" DeepL API integration"""
auth_key = "6bd0d1df-69ed-1ff3-**************"  # (auth_key) Replace with your key
translator = deepl.Translator(auth_key)

# DeepL API testing
result = translator.translate_text("Hello, world!", target_lang="FR")
# print(result.text)  # "Bonjour, le monde !"

do_translation = True

# all_pages = [26, 27, 28, 29, 30, 39, 50, 52, 53, 110, 111, 112, 116, 196, 206, 262, 265, 275, 276, 314, 317, 318, 319, 320, 321, 322, 324, 325, 326, 327, 328, 329, 401, 405, 406, 407, 409, 410, 411, 412, 413, 415, 418, 419, 422, 423, 424, 425, 426, 427, 428, 431, 432, 433, 434, 437, 438, 439, 440, 441, 442, 443, 444, 445, 462, 463, 468, 469, 470, 471, 472, 475, 476, 477, 478, 483, 484, 491, 492, 493, 502, 506, 520, 542, 544, 549, 568, 574, 578, 581, 582, 601, 605, 611, 626, 627, 631, 649, 654, 655, 658, 659, 660, 662, 664, 665, 666, 667, 670, 671, 672, 674, 678, 681, 682, 683, 684, 686, 688, 689, 690, 691, 693, 695, 696, 697, 698, 699, 700]

all_pages = [26, 27, 28, 29, 30, 39, 50, 52, 53, 110]

def buildPresentation():
    prs = Presentation()
    presentationLength = len(all_pages)
    pictureSlide = 0

    if (slideOption == 2):
        pictureSlide = 1
    # web scrapes the URL to get all needed information from the page
    for i in range(0, presentationLength):
        URL = "https://projectabcd.com/display_the_dress.php?id=" + str(all_pages[i])
        page = requests.get(URL, headers={"User-Agent": "html"})
        soup = BeautifulSoup(page.content, "html.parser")
        logo = soup.find("img")
        printLogo = logo.attrs["src"]
        logoURL = "http://projectabcd.com/" + printLogo
        r = requests.get(logoURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
            with open(basename(printLogo), "wb") as f:
                r.raw.decode_content = True
                shutil.copyfileobj(r.raw, f)
        pageInfo = soup.find("div", class_="containerTitle")
        pageInfoImg = soup.find("div", class_="container")
        name = pageInfo.find("h2", class_="headTwo")
        printName = name.text
        image = pageInfoImg.find("div", class_="containerImage")
        img = image.find("image", class_="image")
        printImage = img.get("src")
        pictureURL = "http://projectabcd.com/" + printImage
        r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
            with open(basename(printImage), "wb") as f:
                r.raw.decode_content = True
                shutil.copyfileobj(r.raw, f)

        pagetext = pageInfoImg.find("div", class_="containerText")
        description = pagetext.find("p", class_="words")
        printDescription = (description.text)
        fact = description.find_next_sibling("p")
        printFact = fact.text

        # We now got printDescription and printFact
        if do_translation:
            if inp == "DeepL":  # DeepL API
                print('DeepL API Engine is translating the text...')
                printDescription = DeepL(printDescription)
                printFact = DeepL(printFact)

            elif inp == 'Azure': #Azure API
                print('MS-Azure API Engine is translating the text...')
                printDescription = translate_text_Azure(printDescription, "en", dest_language)
                printFact = translate_text_Azure(printFact, "en", dest_language)

            elif inp == "Google":  # GoogleTranslate API
                print('GoogleTranslate API Engine is translating the text...')
                printDescription = GoogleTranslator1(printDescription, dest_language)
                printFact = GoogleTranslator2(printFact, dest_language)

            elif inp == 'ChatGPT':  # ChatGPT API

                print('ChatGPT API Engine is translating the text...')
                printDescription = translate_text_CHATGPTAPI(printDescription, 'en', dest_language)
                printFact = translate_text_CHATGPTAPI(printDescription, 'en', dest_language)

        else:
            print("Translations failed")

        # creates the slide presentation if slide option 1 is choosen
        if (slideOption == 1):
            # creates the slides and sets layout preferences
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            prs.slide_width = pptx.util.Inches(8)
            prs.slide_height = pptx.util.Inches(11)
            # places the logo on the slide
            logoHolder = slide.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(0),
                                                  width=pptx.util.Inches(1), height=pptx.util.Inches(1))
            # places the title on the slide
            titleBox = slide.shapes.add_textbox(pptx.util.Inches(2.5), pptx.util.Inches(.5),
                                                width=pptx.util.Inches(3), height=pptx.util.Inches(1))
            titleBoxtf = titleBox.text_frame
            title = titleBoxtf.add_paragraph()
            title.text = printName
            title.font.name = titleFont
            title.font.size = Pt(titleSize)
            # places the picture on the slide
            pictureHolder = prs.slides[i].shapes
            pictureHolder.add_picture(basename(printImage), pptx.util.Inches(2.5), pptx.util.Inches(2),
                                      width=pptx.util.Inches(3), height=pptx.util.Inches(4))
            # creates a textbox for the description and fun fact
            contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6),
                                                  width=pptx.util.Inches(6), height=pptx.util.Inches(5))
            contentBoxtf = contentBox.text_frame
            contentBoxtf.word_wrap = True
            descriptionTitle = contentBoxtf.add_paragraph()
            descriptionTitle.font.name = textFont
            descriptionTitle.font.bold = True
            descriptionTitle.font.size = Pt(textSize)
            descriptionTitle.text = "Description: "
            descriptionParagraph = contentBoxtf.add_paragraph()
            descriptionParagraph.font.name = textFont
            descriptionParagraph.font.size = Pt(textSize)
            descriptionParagraph.text = printDescription
            FunFactTitle = contentBoxtf.add_paragraph()
            FunFactTitle.font.name = textFont
            FunFactTitle.font.bold = True
            FunFactTitle.font.size = Pt(textSize)
            FunFactTitle.text = "\nFun Fact:"
            FunFactParagraph = contentBoxtf.add_paragraph()
            FunFactParagraph.font.name = textFont
            FunFactParagraph.font.size = Pt(textSize)
            FunFactParagraph.text = printFact
        # creates the slide presentation if slide option 4 is chosen
        elif (slideOption == 4):
             # creates slide preferences
            slide_layout = prs.slide_layouts[6]
            prs.slide_width = pptx.util.Inches(11)
            prs.slide_height = pptx.util.Inches(8)
            slide2 = prs.slides.add_slide(slide_layout)
            # places picture to cover whole slide
            pictureHolder = prs.slides[pictureSlide].shapes
            pictureHolder.add_picture(basename(printImage), pptx.util.Inches(6), pptx.util.Inches(1),
                                      width=pptx.util.Inches(4), height=pptx.util.Inches(6))
            # creates next slide

            # place logo on the slide
            logoHolder = slide2.shapes.add_picture(basename(printLogo), pptx.util.Inches(5.5), pptx.util.Inches(0),
                                                   width=pptx.util.Inches(1), height=pptx.util.Inches(1))
            # places the title
            titleBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(1.5),
                                                 width=pptx.util.Inches(2), height=pptx.util.Inches(1))
            titleBoxtf = titleBox.text_frame
            title = titleBoxtf.add_paragraph()
            title.text = printName
            title.font.size = Pt(titleSize)
            title.font.name = titleFont
            # creates textbox for description and fun fact
            contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(1),
                                                   width=pptx.util.Inches(3), height=pptx.util.Inches(4))
            contentBoxtf = contentBox.text_frame
            contentBoxtf.word_wrap = True
            descriptionTitle = contentBoxtf.add_paragraph()
            descriptionTitle.font.name = textFont
            descriptionTitle.font.bold = True
            descriptionTitle.font.size = Pt(textSize)
            descriptionTitle.text = "Description: "
            descriptionParagraph = contentBoxtf.add_paragraph()
            descriptionParagraph.font.name = textFont
            descriptionParagraph.font.size = Pt(textSize)
            descriptionParagraph.text = printDescription
            FunFactTitle = contentBoxtf.add_paragraph()
            FunFactTitle.font.bold = True
            FunFactTitle.font.name = textFont
            FunFactTitle.font.size = Pt(textSize)
            FunFactTitle.text = "\nFun Fact:"
            FunFactParagraph = contentBoxtf.add_paragraph()
            FunFactParagraph.font.name = textFont
            FunFactParagraph.font.size = Pt(textSize)
            FunFactParagraph.text = printFact
            pictureSlide = pictureSlide + 1

        # creates the slide presentation if slide option 5 is chosen
        elif (slideOption == 5):
            # creates slide preferences
            slide_layout = prs.slide_layouts[6]
            prs.slide_width = pptx.util.Inches(11)
            prs.slide_height = pptx.util.Inches(8)
            slide2 = prs.slides.add_slide(slide_layout)
            # places picture to cover whole slide
            pictureHolder = prs.slides[pictureSlide].shapes
            pictureHolder.add_picture(basename(printImage), pptx.util.Inches(1), pptx.util.Inches(1),
                                      width=pptx.util.Inches(4), height=pptx.util.Inches(6))
            # creates next slide

            # place logo on the slide
            logoHolder = slide2.shapes.add_picture(basename(printLogo), pptx.util.Inches(5), pptx.util.Inches(0),
                                                   width=pptx.util.Inches(1), height=pptx.util.Inches(1))
            # places the title
            titleBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(1.5),
                                                 width=pptx.util.Inches(2), height=pptx.util.Inches(1))
            titleBoxtf = titleBox.text_frame
            title = titleBoxtf.add_paragraph()
            title.text = printName
            title.font.size = Pt(titleSize)
            title.font.name = titleFont
            # creates textbox for description and fun fact
            contentBox = slide2.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(1),
                                                   width=pptx.util.Inches(3), height=pptx.util.Inches(4))
            contentBoxtf = contentBox.text_frame
            contentBoxtf.word_wrap = True
            descriptionTitle = contentBoxtf.add_paragraph()
            descriptionTitle.font.name = textFont
            descriptionTitle.font.bold = True
            descriptionTitle.font.size = Pt(textSize)
            descriptionTitle.text = "Description: "
            descriptionParagraph = contentBoxtf.add_paragraph()
            descriptionParagraph.font.name = textFont
            descriptionParagraph.font.size = Pt(textSize)
            descriptionParagraph.text = printDescription
            FunFactTitle = contentBoxtf.add_paragraph()
            FunFactTitle.font.bold = True
            FunFactTitle.font.name = textFont
            FunFactTitle.font.size = Pt(textSize)
            FunFactTitle.text = "\nFun Fact:"
            FunFactParagraph = contentBoxtf.add_paragraph()
            FunFactParagraph.font.name = textFont
            FunFactParagraph.font.size = Pt(textSize)
            FunFactParagraph.text = printFact
            pictureSlide = pictureSlide + 1

    prs.save(output_file)
    return output_file



# MS-Azure API
def translate_text_Azure(text, source_language, target_language):
    azure_subscription_key = "9f3df45b03754a2a827bb889aa2d571d"
    azure_endpoint = "https://api.cognitive.microsofttranslator.com/"

    headers = {
        "Ocp-Apim-Subscription-Key": azure_subscription_key,
        "Content-type": "application/json",
        "Ocp-Apim-Subscription-Region": "centralus"
    }

    url = f"{azure_endpoint}/translate?api-version=3.0&from={source_language}&to={target_language}"

    body = [{"text": text}]
    response = requests.post(url, headers=headers, json=body)

    if response.status_code == 200:
        translated_text = response.json()[0]["translations"][0]["text"]
        return translated_text
    else:
        print(f"Translation error: {response.text}")
        return None


# GoogleTranslateAPI
def GoogleTranslator1(printDescription, dest_language):
    t1 = GoogleTranslator(source='auto', target=dest_language).translate(printDescription)

    return t1


def GoogleTranslator2(printFact, dest_language):
    t2 = GoogleTranslator(source='auto', target=dest_language).translate(printFact)

    return t2


# Function to translate text using ChatGPT API
APIKey = 'sk-ByZkXS3Ur7YQHTCVaUVbT3****************' # replace with your key

openai.api_key = APIKey


def translate_text_CHATGPTAPI(text, source_language, target_language):
    prompt = f"Translate the following '{source_language}' text to '{target_language}': {text}"

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that translates text."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=150,
        n=1,
        stop=None,
        temperature=0.5,
    )

    translation = response.choices[0].message.content.strip()
    return translation

#DeepL API
def DeepL(printDescription):
    result = (
        translator.translate_text(printDescription, target_lang=dest_language_DeepL, split_sentences=1)).__str__().encode(
        "utf-8")

    return result


ppt_file = buildPresentation()
os.startfile(ppt_file)
