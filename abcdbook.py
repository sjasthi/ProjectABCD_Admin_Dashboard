import random
import sys, io, requests, threading, webbrowser, urllib, os, platform, openpyxl, string, textwrap, re, time, textstat
from xml.dom.expatbuilder import FragmentBuilderNS
import tkinter as tk
from tkinter import ttk, messagebox
from pptx import Presentation
import pptx.util
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import pandas as pd
from pandas import Series, DataFrame
from concurrent.futures import ThreadPoolExecutor, as_completed
from textblob import TextBlob
from bs4 import BeautifulSoup
import wikipediaapi
import openai

# pip install googletrans
# pip install googletrans==4.0.0-rc1
import googletrans

openai.api_key = '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'

ROOT_WIDTH = 1000 # app window width
ROOT_HEIGHT = 600 # app window height

root = tk.Tk()
root.title("Project ABCD Admin Panel")
sw_placement = int(root.winfo_screenwidth()/2 - ROOT_WIDTH/2) # to place at half width of screen
sh_placement = int(root.winfo_screenheight()/2 - ROOT_HEIGHT/2) # to place at half height of screen
root.geometry(f"{ROOT_WIDTH}x{ROOT_HEIGHT}+{sw_placement}+{sh_placement}")
root.minsize(ROOT_WIDTH, ROOT_HEIGHT)

# Set sys.stdout to use utf-8 encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

MAIN_FONT = ("helvetica", 12)
LABEL_FONT = ("helvetica bold", 14)
LANGUAGES = [
    "Telugu",
    "Hindi",
    "Spanish"
]

# creates preferences dictionary from preferences.txt
# creates default preferences.txt file if one does not exist
try:
    with open("preferences.txt", "r", encoding="utf8") as file:
        lines = file.readlines()
        preferences = {}

        for line in lines:
            key, value = line.split('=')
            preferences[key.strip()] = value.strip().replace('“', '').replace('”', '').replace('"', '').replace("'", '')
except FileNotFoundError:
    preferences = {
        "TEXT_SIZE" : "14",
        "TEXT_FONT" : "Times New Roman",
        "TITLE_SIZE" : "32",
        "TITLE_FONT" : "Arial",
        "SUBTITLE_SIZE" : "24",
        "SUBTITLE_FONT" : "Arial",
        "PIC_WIDTH" : "720",
        "PIC_HEIGHT" : "1040",
        "PUZ_WIDTH" : "20",
        "PUZ_HEIGHT" : "20",
        "WORD_COUNT": "10"
    }
    tk.messagebox.showwarning(title='Warning', message='No preferences.txt file exists in directory. Default preferences.txt will be created and used.')
    print('No preferences.txt file exists in directory. Default preferences.txt will be created and used.')
    with open('preferences.txt', 'w') as f:
        f.write('TEXT_SIZE = 14\nTEXT_FONT = Times New Roman\nTITLE_SIZE = 32\nTITLE_FONT = Arial\nSUBTITLE_SIZE = 24\nSUBTITLE_FONT = Arial\nPIC_WIDTH = 720\nPIC_HEIGHT = 1040')

'''
Gathers data from API
'''
def downloadAPIData(url, id_number):
    try:
        response = requests.get(url, headers={"User-Agent": "XY"})
        # append dress info to dress data if response status_code == 200
        if response.ok:
            return response.json()['data']
        else:
            print(f'Request for dress ID: {id_number} failed.')
    except requests.exceptions.RequestException as e:
        print(f'-- DEBUG -- in downloadAPIData: {e}')
    except Exception as e:
        # tk.messagebox.showerror(title="Error", message=f'Could not make connection!\n\nError: {e}')
        print(f'Error: {e}')
    
'''
Sets up and starts threads for gathering API data
'''
def apiRunner():
    # dress ids in entry field
    dress_ids = getSlideNumbers()

    # create list of all urls to send requests to
    url_list = []
    for id_number in dress_ids:
        url_list.append(f'https://abcd2.projectabcd.com/api/getinfo.php?id={id_number}')

    dress_data = [] # dress data from API
    threads= [] # working threads

    # create progress bar
    progress_window, pb, percent_label = progress_bar('Retrieving API Data')

    # spins up 10 threads at a time and stores retrieved data into dress_data upon completion
    with ThreadPoolExecutor(max_workers=10) as exec:
        for index, url in enumerate(url_list):
            threads.append(exec.submit(downloadAPIData, url, dress_ids[index]))
        
        complete = 0 # number of threads that have finished
        for task in as_completed(threads):
            complete += 1
            pb['value'] = (complete/len(dress_ids))*100 # calculate percentage of data retrieved
            percent_label.config(text=f'Retrieving API Data...{int(pb["value"])}%') # update completion percent label

            if task.result() is not None:
                dress_data.append(task.result()) # append retrieved data to dress_data

    progress_window.destroy()
    return dress_data

'''
Gets dress IDs from entry field
'''
def getSlideNumbers():
    # get update for dress number input
    update_dress_list = []
    # get dress numbers from text field
    get_text_field = text_field.get("1.0", "end-1c").split(',')
    
    # add to list
    for number in get_text_field:
        if (number.strip().isnumeric()):
            update_dress_list.append(int(number.strip()))
    
    # remove duplicates
    dress_ids = []
    [dress_ids.append(x) for x in update_dress_list if x not in dress_ids]

    return dress_ids

'''
Downloads dress images
'''
def downloadImages(folder_name, url, img_name):
    try:
        # downloads dress image
        img_url = url
        img_path = f'./{folder_name}/{img_name}'
        opener = urllib.request.build_opener()
        opener.addheaders=[('User-Agent', 'XY')]
        urllib.request.install_opener(opener)
        urllib.request.urlretrieve(img_url, img_path)
    except Exception as e:
        print(f'Error downloading images: {e}')

'''
Sets up and starts threads for downloading dress images
'''
def imageRunner(dress_data):
    # list of all urls to send requests to
    url_list = []
    # list of all image names
    img_name_list = []

    for data in dress_data:
        url_list.append(f'http://projectabcd.com/images/dress_images/{data["image_url"]}')
        img_name_list.append(f'{data["image_url"]}')

    # create progress bar
    progress_window, pb, percent_label = progress_bar('Downloading Images')

    threads= [] # working threads

    # spins up 10 threads at a time and calls downloadImages with url and image name
    with ThreadPoolExecutor(max_workers=10) as exec:
        for index, url in enumerate(url_list):
            threads.append(exec.submit(downloadImages, "images", url, img_name_list[index]))

        complete = 0 # number of threads that have finished
        for task in as_completed(threads):
            complete += 1
            pb['value'] = (complete/len(url_list))*100 # calculate percentage of images downloaded
            percent_label.config(text=f'Downloading Images...{int(pb["value"])}%') # update completion percent label
    
    progress_window.destroy() # close progress bar window

'''
Generates progress bar
'''
def progress_bar(title):
    # progress bar window for API data retrieval
    progress_window = tk.Toplevel(root)
    progress_window.title(title)
    sw = int(progress_window.winfo_screenwidth()/2 - 450/2)
    sh = int(progress_window.winfo_screenheight()/2 - 70/2)
    progress_window.geometry(f'450x70+{sw}+{sh}')
    progress_window.resizable(False, False)
    progress_window.attributes('-disable', True)
    progress_window.focus()

    # progress bar custom style
    pb_style = ttk.Style()
    pb_style.theme_use('clam')
    pb_style.configure('green.Horizontal.TProgressbar', foreground='#1ec000', background='#1ec000')

    # frame to hold progress bar
    pb_frame = tk.Frame(progress_window)
    pb_frame.pack()

    # progress bar
    pb = ttk.Progressbar(pb_frame, length=400, style='green.Horizontal.TProgressbar', mode='determinate', maximum=100, value=0)
    pb.pack(pady=10)

    # label for percent complete
    percent_label = tk.Label(pb_frame, text=f'{title}...0%')
    percent_label.pack()

    return progress_window, pb, percent_label

'''
Translate text to selected language 
'''
def translateText(text):
    # Build the translator
    if translate.get() == 1:
        if language.get() == "Telugu":
            dest_language = 'te'
        elif language.get() == "Hindi":
            dest_language = 'hi'
        elif language.get() == "Spanish":
            dest_language = 'es'
    else:
        return text # return text if english
    try:
        # Create the translator
        translator = googletrans.Translator()
        translated_text = translator.translate(text, dest = dest_language)
        return translated_text.text
    except Exception as e:
        print("Error:", e)
        return ""

'''
Sorts dress data based on user selection
'''
def sortDresses(dress_data):
    if sort_order.get() == 1:
        return sorted(dress_data, key=lambda x : str(x['name']).lower())
    elif sort_order.get() == 2:
        return sorted(dress_data, key=lambda x : x['id'])
    elif sort_order.get() == 3:
        return dress_data

'''
Opens file depending on OS
'''
def openFile(file_name):
    current_os = platform.system()
    try:
        if current_os == "Windows":
            print(f"-- DEBUG -- Windows {file_name}")
            os.system(f"start {file_name}")
        elif current_os == "Darwin":
            print(f"-- DEBUG -- Darwin {file_name}")
            os.system(f"open {file_name}")
        elif current_os == "Linux":
            print(f"-- DEBUG -- Linux {file_name}")
            os.system(f"xdg-open {file_name}")
        else:
            print("Error: Cannot open file " + current_os + " not supported.")
    except Exception as e:
        print("Error:", e)

'''
Create the Dress Name Title textbox
'''
def add_title_box(slide, dress_name, left, top, width, height):
    title = slide.shapes.title
    title.left = Inches(left)
    title.top = Inches(top)
    title.width = Inches(width)
    title.height = Inches(height)
    title.text = f'{dress_name.upper()}'

    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x09, 0x09, 0x82)
    title.text_frame.paragraphs[0].text = title.text_frame.paragraphs[0].text.upper()
    title.text_frame.paragraphs[0].font.name = title_font_var.get()
    title.text_frame.paragraphs[0].font.size = Pt(int(title_size_var.get()))
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

'''
Create the subtitle highlight
'''
def add_subtitle_highlight(slide, left, top, width, height):
    # create the shape
    rectangle1 = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    rectangle1.rotation = 357

    # color fill
    fill1 = rectangle1.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(202, 246, 189)

    # no outline
    line1 = rectangle1.line
    line1.color.rgb = RGBColor(255, 255, 255)

    # no shaddow
    shadow1 = rectangle1.shadow
    shadow1.inherit = False

'''
Create the Description subtitle
'''
def add_description_subtitle(slide, left, top, width, height):
    description_subtitle_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    description_subtitle_frame = description_subtitle_box.text_frame

    description_subtitle = description_subtitle_frame.add_paragraph()
    description_subtitle.font.color.rgb = RGBColor(0x6E, 0xD8, 0xFF)
    description_subtitle.font.bold = True
    description_subtitle.font.name = subtitle_font_var.get()
    description_subtitle.text = translateText("DESCRIPTION:")

    # make text smaller for layout 4
    if layout.get() == 4:
        description_subtitle.font.size = Pt(16)
    else:
        description_subtitle.font.size = Pt(int(subtitle_size_var.get()))

'''
Add the Description text
'''
def add_description_text(slide, dress_description, left, top, width, height):
    # description - text (left, top, width, height)
    description_text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    description_text_frame = description_text_box.text_frame

    description_text_frame.word_wrap = True
    description_text = description_text_frame.add_paragraph()
    description_text.font.name = text_font_var.get()
    description_text.text = f'{translateText(dress_description)}'

    # make text smaller for layout 4
    if layout.get() == 4:
        description_text.font.size = Pt(11)
    else:
        description_text.font.size = Pt(int(text_size_var.get()))

'''
Create Did You Know subtitle
'''
def add_did_you_know_subtitle(slide, left, top, width, height):
    did_you_know_subtitle_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    did_you_know_subtitle_frame = did_you_know_subtitle_box.text_frame

    did_you_know_subtitle = did_you_know_subtitle_frame.add_paragraph()
    did_you_know_subtitle.font.color.rgb = RGBColor(0x6E, 0xD8, 0xFF)
    did_you_know_subtitle.font.bold = True
    did_you_know_subtitle.font.name = subtitle_font_var.get()
    did_you_know_subtitle.text = translateText("DID YOU KNOW?")

    # make text smaller for layout 4
    if layout.get() == 4:
        did_you_know_subtitle.font.size = Pt(16)
    else:
        did_you_know_subtitle.font.size = Pt(int(subtitle_size_var.get()))

'''
Create Did You Know text
'''
def add_did_you_know_text(slide, dress_did_you_know, left, top, width, height):
    did_you_know_text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    did_you_know_text_frame = did_you_know_text_box.text_frame

    did_you_know_text_frame.word_wrap = True
    did_you_know_text = did_you_know_text_frame.add_paragraph()
    did_you_know_text.font.name = text_font_var.get()
    did_you_know_text.text = f'{translateText(dress_did_you_know)}'

    # make text smaller for layout 4
    if layout.get() == 4:
        did_you_know_text.font.size = Pt(11)
    else:
        did_you_know_text.font.size = Pt(int(text_size_var.get()))

'''
Add the dress image 
'''
def add_image(slide, dress_info, left, top):
    image_width_px = int(pic_width_var.get())
    image_height_px = int(pic_height_var.get())
    image_width_inch = image_width_px / 96
    image_height_inch = image_height_px / 96

    # change image size for layout 4 
    if layout.get() == 4:
        image_width_inch = 3.71
        image_height_inch = 4.94

    try:
        picture = slide.shapes.add_picture(f'./images/Slide{dress_info["id"]}.png', 0, Inches(0.83), Inches(image_width_inch), Inches(image_height_inch))
    except FileNotFoundError:
        # tk.messagebox.showerror(title="Error", message="When the Download Images check box is not checked make sure you have an images \
        #                                                 directory in the root of this project that includes the correct images. (example: 1 == Slide1.PNG)")
        print(f'Image Slide{dress_info["id"]}.png Not Found!')

'''
Create the Dress Id & Page Number text
'''
def add_numbering(slide, dress_info, index, left, top, width, height, left2, top2, width2, height2):
    # show both page num & dress id
    if numbering.get() == 1:
        # dress id (left, top, width, height)
        dress_id_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        dress_id_box_frame = dress_id_box.text_frame

        dress_id = dress_id_box_frame.add_paragraph()
        dress_id.font.name = text_font_var.get()
        dress_id.text = translateText(f'Dress ID: {dress_info["id"]}')
        
        # aligns text left & make text smaller
        if layout.get() == 4:
            dress_id.alignment = PP_ALIGN.LEFT
            dress_id.font.size = Pt(11)
        else:
            dress_id.alignment = PP_ALIGN.RIGHT
            dress_id.font.size = Pt(int(text_size_var.get()))

        # page number (left, top, width, height)
        page_number_box = slide.shapes.add_textbox(Inches(left2), Inches(top2), Inches(width2), Inches(height2))
        page_number_box_frame = page_number_box.text_frame

        page_number = page_number_box_frame.add_paragraph()
        page_number.font.name = text_font_var.get()
        page_number.text = translateText(f'Page No. {index+1}')

        # aligns text left & make text smaller
        if layout.get() == 4:
            page_number.alignment = PP_ALIGN.LEFT
            page_number.font.size = Pt(11)
        else:
            page_number.alignment = PP_ALIGN.RIGHT
            page_number.font.size = Pt(int(text_size_var.get()))
    # show page num or dress id
    elif numbering.get() == 2 or numbering.get() == 3:
        number_box = slide.shapes.add_textbox(Inches(left2), Inches(top2), Inches(width2), Inches(height2))
        number_box_frame = number_box.text_frame

        number_box = number_box_frame.add_paragraph()
        number_box.font.name = text_font_var.get()

        # aligns text left & make text smaller
        if layout.get() == 4:
            number_box.alignment = PP_ALIGN.LEFT
            number_box.font.size = Pt(11)
        else:
            number_box.alignment = PP_ALIGN.RIGHT
            number_box.font.size = Pt(int(text_size_var.get()))

        # show page num
        if numbering.get() == 2:
            number_box.text = translateText(f'Page No. {index+1}')
        # show dress id
        elif numbering.get() == 3:
            number_box.text = translateText(f'Dress ID: {dress_info["id"]}')

'''
Closes the pop alert message
'''
def close_popup(popup):
    popup.destroy()

'''
Updates the timer for the alert popup
'''
def update_timer(popup, timer_label, seconds_left):
    timer_label.config(text=f"Auto closes in {seconds_left} sec.")

    if seconds_left > 0:
        # Update the timer every second
        root.after(1000, update_timer, popup, timer_label, seconds_left - 1)
    else:
        close_popup(popup)

'''
Displays the alert message
'''
def show_error_popup(text_message, duration_num):
    popup = tk.Toplevel(root)
    popup.title("Alert Message")
    
    label = tk.Label(popup, text=text_message)
    label.pack(padx=10, pady=10)

    timer_label = tk.Label(popup, text="")
    timer_label.pack(pady=5)

    duration = duration_num
    update_timer(popup, timer_label, duration)

'''
Performs update once generate button clicked
'''
def generateBook():
    # check if Generate from Local is active
    if gen_local.get() == 1:
        # get update for dress number input
        update_dress_list = []
        # get dress numbers from text field
        get_text_field = text_field.get("1.0", "end-1c").split(',')
        
        # add to list
        for number in get_text_field:
            if (number.strip().isnumeric()):
                update_dress_list.append(int(number.strip()))
        
        # remove duplicates
        dress_ids = []
        [dress_ids.append(x) for x in update_dress_list if x not in dress_ids]

        # path to local Excel data
        file_path = "APIData.xlsx"

        # gets and cleans dress data from Excel file
        sheet_dress_data = pd.read_excel(file_path)
        sheet_dress_data.dropna(subset=['id'], inplace=True) # drops any rows with na ID
        sheet_dress_data['description'].fillna('', inplace=True) # removes na/nan from description column
        sheet_dress_data['did_you_know'].fillna('', inplace=True) # removes na/nan from did_you_know column
        sheet_dress_data['description'] = sheet_dress_data['description'].astype(str).apply(openpyxl.utils.escape.unescape) # convert escaped strings to ASCII
        sheet_dress_data['did_you_know'] = sheet_dress_data['did_you_know'].astype(str).apply(openpyxl.utils.escape.unescape) # Convert escaped strings to ASCII
        
        # holds dress data from local Excel sheet
        dress_data = [] 
        # cycle through dress_ids and append data from excel sheet to dress_data
        for id in dress_ids:
            row = sheet_dress_data.loc[id-1]
            dress_data.append({'id':row.loc['id'], 'name':row.loc['name'], 'description':row.loc['description'], 'did_you_know':row.loc['did_you_know']})   
    else:
        dress_data = apiRunner() # gather all dress data from api
    
    # sort dress data
    sorted_dress_data = sortDresses(dress_data)

    # if there is no image in the local folder then download from web
    if download_imgs.get() == 0:
        if not os.path.exists('./images'):
            os.makedirs('./images')
        if not os.listdir('./images'):
            text_message = "Local folder is empty. Attempting to grab image(s) from web."
            duration_num = 4
            show_error_popup(text_message, duration_num)
            time.sleep(duration_num+1)
            imageRunner(sorted_dress_data)

    # download images from web if download images check box is selected
    if download_imgs.get() == 1:
        # creates directory to save images if one does not exist
        if not os.path.exists('./images'):
            os.makedirs('./images')
        imageRunner(sorted_dress_data) # download images for each dress in list

    # create powerpoint
    prs = Presentation() # create the pptx presentation
    ppt_file_name = "abcdbook.pptx"
    file_name = "abcdbook.pptx"
    count = 0
    while os.path.exists(file_name): # check if file name exist,
        count += 1
        file_name = f"{os.path.splitext(ppt_file_name)[0]}({count}).pptx" # if file name exisit create new filename

    # progress bar window for API data retrieval
    progress_window = tk.Toplevel(root)
    progress_window.title('Creating Book')
    sw = int(progress_window.winfo_screenwidth()/2 - 450/2)
    sh = int(progress_window.winfo_screenheight()/2 - 70/2)
    progress_window.geometry(f'450x70+{sw}+{sh}')
    progress_window.resizable(False, False)
    progress_window.attributes('-disable', True)
    progress_window.focus()

    # progress bar custom style
    pb_style = ttk.Style()
    pb_style.theme_use('clam')
    pb_style.configure('green.Horizontal.TProgressbar', foreground='#1ec000', background='#1ec000')

    # frame to hold progress bar
    pb_frame = tk.Frame(progress_window)
    pb_frame.pack()

    # progress bar
    pb = ttk.Progressbar(pb_frame, length=400, style='green.Horizontal.TProgressbar', mode='determinate', maximum=100, value=0)
    pb.pack(pady=10)

    # label for percent complete
    percent_label = tk.Label(pb_frame, text='Creating Book...0%')
    percent_label.pack()
    complete = 0

    # get dress info for items in list & translate
    for index, dress_info in enumerate(sorted_dress_data):
        left = None
        image_left = None
        
        dress_name = dress_info['name']
        dress_description = dress_info['description']
        dress_did_you_know = dress_info['did_you_know']
        dress_description_len = len(dress_description)

        #--------------------------------Portrait--------------------------------
        # PORTRAIT MODE
        if layout.get() == 1 or layout.get() == 4: 
            prs.slide_width = pptx.util.Inches(7.5) # define slide width
            prs.slide_height = pptx.util.Inches(10.83) # define slide height
            slide_layout = prs.slide_layouts[5] # use slide with only title
            slide_layout2 = prs.slide_layouts[6] # use empty slide

            # LAYOUT 1 == picture on left page - text on right page - two page
            if layout.get() == 1:
                slide_empty = prs.slides.add_slide(slide_layout2) 
                slide_title = prs.slides.add_slide(slide_layout) 

                add_image(slide_empty, dress_info, 0, 0)
                add_title_box(slide_title, dress_name, 0, 0.15, 7.5, 0.91) 
                add_subtitle_highlight(slide_title, 0.37, 1.58, 2.44, 0.3) # description - highlight box
                add_description_subtitle(slide_title, 0.28, 1.07, 6.94, 0.51)
                add_description_text(slide_title, dress_description, 0.28, 1.65, 6.94, 5.99)
                add_subtitle_highlight(slide_title, 0.37, 8.36, 2.78, 0.3) # did you know - highlight box
                add_did_you_know_subtitle(slide_title, 0.28, 7.87, 6.94, 0.51)
                add_did_you_know_text(slide_title, dress_did_you_know, 0.28, 8.46, 6.94, 1.04)
                add_numbering(slide_title, dress_info, index, 4.47, 10.06, 1.28, 0.34, 5.94, 10.06, 1.28, 0.34)

            # LAYOUT 4 == picture on left - text on right - single page
            elif layout.get() == 4: 
                slide_title = prs.slides.add_slide(slide_layout) 

                add_image(slide_title, dress_info, 0, 1.39)
                add_title_box(slide_title, dress_name, 0, 0.09, 7.5, 0.91) 
                add_subtitle_highlight(slide_title, 3.71, 1.21, 1.83, 0.23) # description - highlight box
                add_description_subtitle(slide_title, 3.63, 0.88, 3.71, 0.37)
                add_description_text(slide_title, dress_description, 3.63, 1.35, 3.71, 7.57)

                # adjust the text height based on text length
                if dress_description_len < 600:
                    add_subtitle_highlight(slide_title, 3.72, 5.83, 1.83, 0.23) # did you know - highlight box
                    add_did_you_know_subtitle(slide_title, 3.63, 5.42, 3.71, 0.37)
                    add_did_you_know_text(slide_title, dress_did_you_know, 3.63, 5.94, 3.71, 0.91)

                elif dress_description_len > 600 and dress_description_len < 1300:
                    add_subtitle_highlight(slide_title, 3.72, 7.99, 1.83, 0.23) # did you know - highlight box
                    add_did_you_know_subtitle(slide_title, 3.63, 7.58, 3.71, 0.37)
                    add_did_you_know_text(slide_title, dress_did_you_know, 3.63, 8.1, 3.71, 0.91)
                    
                else:
                    add_subtitle_highlight(slide_title, 3.72, 9.32, 1.83, 0.23) # did you know - highlight box
                    add_did_you_know_subtitle(slide_title, 3.63, 8.91, 3.71, 0.37)
                    add_did_you_know_text(slide_title, dress_did_you_know, 3.63, 9.43, 3.71, 0.91)

                add_numbering(slide_title, dress_info, index, 0.49, 6.46, 1.33, 0.27, 1.95, 6.46, 1.33, 0.27)
 
        #--------------------------------Landscape--------------------------------
        # LANDSCAPE MODE
        elif layout.get() == 2 or layout.get() == 3:
            left = None
            image_left = None
            rectangle_left = None

            # slide size (left, top, width, height)
            prs.slide_width = pptx.util.Inches(16.18) # define slide width
            prs.slide_height = pptx.util.Inches(12.53) # define slide height
            slide_layout = prs.slide_layouts[5] # use empty slide
            slide = prs.slides.add_slide(slide_layout) # add empty slide to pptx

            # LAYOUT 2 == picture on right - text on left
            if layout.get() == 2:
                rectangle_left = 0.4
                left = 0.34
                image_left = 8.45
                image_top = 1.17
                numbering1_left = 1.05
                numbering2_left = 2.53

            # LAYOUT 3 == picture on left - text on right
            elif layout.get() == 3:
                image_left = 0.25
                image_top = 1.17
                rectangle_left = 8.15
                left = 8.09
                numbering1_left = 12.78
                numbering2_left = 14.26

            add_title_box(slide, dress_name, 0, 0.15, 16.18, 0.91) 
            add_subtitle_highlight(slide, rectangle_left, 1.75, 2.44, 0.3) # decription - highlight box
            add_description_subtitle(slide, left, 1.26, 7.81, 0.51)
            add_description_text(slide, dress_description, left, 1.88, 7.81, 5.35)
            add_subtitle_highlight(slide, rectangle_left, 8.04, 2.76, 0.28) # did you know - highlight box
            add_did_you_know_subtitle(slide, left, 7.56, 7.81, 0.51)
            add_did_you_know_text(slide, dress_did_you_know, left, 8.19, 7.81, 1.11)
            add_numbering(slide, dress_info, index, numbering1_left, 11.08, 1.28, 0.34, numbering2_left, 11.08, 1.28, 0.34)
            add_image(slide, dress_info, image_left, image_top)
        complete += 1
        pb['value'] = (complete/len(sorted_dress_data))*100 # calculate percentage of images downloaded
        percent_label.config(text=f'Creating Book...{int(pb["value"])}%') # update completion percent label
    try:
        prs.save(file_name)
    except Exception as e:
        print(f"-- DEBUG -- saving presentation: {e}")
    finally:
        book_gen_generate_button.config(state="normal")
    
    progress_window.destroy() # close progress bar window

    openFile(file_name)

'''
Helper function to wrap text
'''
def wrap(string, length=150):
    return '\n'.join(textwrap.wrap(string, length))

'''
Generates treeview table of data
'''
def generate_table(table_data, report_name, column_headers, row_h, col_w, anchor_point, num_buttons=2):
    # create window to display table
    table_window = tk.Toplevel(root)
    table_window.title(report_name)
    table_window.geometry(f"1000x600")
    table_window.minsize(1000,600)

    # create frame to hold table
    table_frame = tk.Frame(table_window)
    table_frame.pack_propagate(False)
    table_frame.place(x=0, y=0, relwidth=1, relheight=.89, anchor="nw")

    # using style to set row height and heading colors
    style = ttk.Style()
    style.theme_use('clam')
    style.configure('Treeview', rowheight=row_h)
    style.configure('Treeview.Heading', background='#848484', foreground='white')

    # vertical scrollbar
    table_scrolly = tk.Scrollbar(table_frame)
    table_scrolly.pack(side="right", fill='y')
    # horizontal scrollbar
    table_scrollx = tk.Scrollbar(table_frame, orient='horizontal')
    table_scrollx.pack(side="bottom", fill='x')

    # use ttk Treeview to create table
    table = ttk.Treeview(table_frame, yscrollcommand=table_scrolly.set, xscrollcommand=table_scrollx.set, columns=column_headers, show='headings')

    # configure the scroll bars with the table
    table_scrolly.config(command=table.yview)
    table_scrollx.config(command=table.xview)

    # create the headers and set column variables
    for index, column_header in enumerate(column_headers):
        table.heading(column_header, text=column_header)
        if index == 0:
            table.column(column_header, width=75, stretch=False)
        elif index == 1:
            table.column(column_header, width=145, stretch=False)
        else:
            table.column(column_header, width=col_w, stretch=False, anchor=anchor_point)

    # pack table into table_frame
    table.pack(fill='both', expand=True)

    # fill table with difference report data
    for index, data in enumerate(table_data):
        # word wrap text
        for i, cell in enumerate(data): 
            if len(str(data[i])) > 2000:
                data[i] = wrap(str(cell), 400)
            else:
                data[i] = wrap(str(cell), 250)

        # if new row, set tag to new
        # if changed row, set tag to changed
        if data[-1] == 'new':
            table.insert(parent='', index=tk.END, values=data, tags=('new',))
        elif data[-1] == 'changed':
            table.insert(parent='', index=tk.END, values=data, tags=('changed',))
        else:
            # if even row, set tag to evenrow
            # if odd row, set tag to oddrow
            if index % 2 == 0:
                table.insert(parent='', index=tk.END, values=data, tags=('evenrow',))
            else:
                table.insert(parent='', index=tk.END, values=data, tags=('oddrow',))

    # color rows
    table.tag_configure('new', background='#BAFFA4')
    table.tag_configure('changed', background='#FFA5A4')
    table.tag_configure('evenrow', background='#e8f3ff')
    table.tag_configure('oddrow', background='#f7f7f7')

    # create button frame and place one table_window
    btn_frame = tk.Frame(table_window)
    btn_frame.pack(side='bottom', pady=15)
    
    if num_buttons == 2:
        # create buttons and pack on button_frame
        btn = tk.Button(btn_frame, text='Export SQL File', font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: exportSQL(table_data, column_headers, report_name))
        btn.pack(side='left', padx=25)

        btn2 = tk.Button(btn_frame, text='Export to HTML', font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: exportHTML(table_data, column_headers, report_name))
        btn2.pack(side='left', padx=25)

    elif num_buttons == 1:
        btn = tk.Button(btn_frame, text='Export to HTML', font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: exportHTML(table_data, column_headers, report_name))
        btn.pack(side='left', padx=25)

'''
Exports table data to Excel file
'''
def exportExcel(data, excel_columns, sheet_name):
    df = pd.DataFrame(data, columns=excel_columns)
    df.to_excel(f'{sheet_name}.xlsx', index=False)

'''
Exports difference report data to SQL update script
'''
def exportSQL(dress_data, column_headers, report_name):
    if report_name == 'difference_report':
        sql_queries = [] # stores sql query

        # cycle through the diff_dress_data and create queries for each
        for index, data in enumerate(dress_data):
            data[1] = str(data[1]).replace('"', '\\"') # replace " with \" so quotations don't mess up query
            data[2] = str(data[2]).replace('"', '\\"') # replace " with \" so quotations don't mess up query
            data[3] = str(data[3]).replace('"', '\\"') # replace " with \" so quotations don't mess up query
            if data[-1] == 'changed':
                sql_queries.append(f'UPDATE dresses\nSET name="{data[1]}", description="{data[2]}", did_you_know="{data[3]}"\nWHERE id={data[0]};\n')
            elif data[-1] == 'new':
                sql_queries.append(f'INSERT INTO dresses (id, name, description, did_you_know)\nVALUES ({data[0]}, "{data[1]}", "{data[2]}", "{data[3]}");\n')
        
        # create path for sql script
        update_script_path = 'abcdbook_SQL_update.sql'
        update_script_name = 'abcdbook_SQL_update'
        count = 1
        while os.path.exists(update_script_path):
            update_script_path = f'{update_script_name}({count}).sql'
            count += 1

        # write sql queries into .sql script
        with open(update_script_path, 'w') as f:
            for query in sql_queries:
                f.write(f'{query}\n')

    elif report_name == 'wiki_link_report':
        with open(f'{report_name}_update.sql', 'w') as sql_file:
            # Create SQL CREATE TABLE statement
            create_table_query = f'CREATE TABLE IF NOT EXISTS resources (\n'
            create_table_query += ', '.join(f'{column} TEXT' for column in column_headers)
            create_table_query += '\n);\n\n'
            sql_file.write(create_table_query)

            # Create SQL INSERT INTO statement
            sql_file.write(f'INSERT INTO resources ({", ".join(column_headers)}) VALUES\n')

            # Iterate through data and write values
            for row in dress_data:
                values = ', '.join(f"'{str(value)}'" for value in row)
                sql_file.write(f'({values}),\n')

            # Remove the trailing comma from the last line
            sql_file.seek(sql_file.tell() - 2)
            sql_file.truncate()

            # Add a semicolon to the end of the SQL script
            sql_file.write(';')

'''
Exports data to JQuery data table HTML page
'''
def exportHTML(data, column_headers, file_name):
    #read in data and create column headers
    table_data = pd.DataFrame(data)
    table_data.columns = column_headers

    #convert table_data to html
    html_table_data = table_data.to_html(table_id='html_table_data', border=0, classes='display')

    #html template to generate JQuery DataTable of data
    html_temp = f"""
    <!DOCTYPE html>
    <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=width-device, initial-scale=1.0">
            <title>{file_name}</title>
            <!--jQuery cdn-->
            <script src="https://code.jquery.com/jquery-3.7.0.js"></script>
            <!--Datatable style-->
            <link rel="stylesheet" href="https:////cdn.datatables.net/1.13.6/css/jquery.dataTables.min.css" />
            <!--Datatable cdn-->
            <script src="https:////cdn.datatables.net/1.13.6/js/jquery.dataTables.min.js"></script>
            <!--Datatable button libraries-->
            <script src="https://cdn.datatables.net/buttons/2.4.2/js/dataTables.buttons.min.js"></script>
            <script src="https://cdn.datatables.net/buttons/2.4.2/js/buttons.html5.min.js"></script>
            <script src="https://cdnjs.cloudflare.com/ajax/libs/jszip/3.10.1/jszip.min.js"></script>
            <script src="https://cdnjs.cloudflare.com/ajax/libs/pdfmake/0.1.53/pdfmake.min.js"></script>
            <script src="https://cdnjs.cloudflare.com/ajax/libs/pdfmake/0.1.53/vfs_fonts.js"></script>
            <!--Initialize datatables-->
            <script>
                $(document).ready(function() {{
                    $('#html_table_data').DataTable({{
                        pageLength: 25,
                        dom: 'Bfrtip',
                        buttons: [
                            'csv', 'excel', 'pdf'
                        ]
                    }},
                    style_table = {{
                        'width': '100%'
                    }});
                }});
            </script>
        </head>
        <body>
            {html_table_data}
        </body>
    </html>
    """

    #write html template into datatable.html file
    with open(f'{file_name}.html', 'w', encoding='utf-8') as f:
        f.write(html_temp)

    #open datatable.html
    webbrowser.open(f'{file_name}.html')

'''
Performs difference report on Excel sheet compared to API data
'''
def diffReport():
    file_path = 'APIData.xlsx' # Change to path where file is located

    dress_ids = sorted(getSlideNumbers()) # gets dress IDs in entry field
    diff_dress_data = [] # data in spreadsheet that is different from API
    api_dress_data = sorted(apiRunner(), key=lambda x : x['id']) # gets dress data from API and sorts by ID

    try:
        # gets and cleans dress data from Excel file
        sheet_dress_data = pd.read_excel(file_path)
        sheet_dress_data.dropna(subset=['id'], inplace=True) # drops any rows with na ID
        sheet_dress_data['description'].fillna('', inplace=True) # removes na/nan from description column
        sheet_dress_data['did_you_know'].fillna('', inplace=True) # removes na/nan from did_you_know column
        sheet_dress_data['description'] = sheet_dress_data['description'].astype(str).apply(openpyxl.utils.escape.unescape) # convert escaped strings to ASCII
        sheet_dress_data['did_you_know'] = sheet_dress_data['did_you_know'].astype(str).apply(openpyxl.utils.escape.unescape) # Convert escaped strings to ASCII

        # cycle through api_dress_data
        for api_data in api_dress_data:
            # row of data with an ID that matches api_data ID
            row = sheet_dress_data.loc[api_data['id']-1] # row of data in spreadsheet

            # check if name, description, or did_you_know is different from the API data
            if row.loc['name'] != api_data['name']:
                new_row = [item for item in row]
                new_row.append('changed')
                diff_dress_data.append(new_row)
                continue
            if row.loc['description'] != api_data['description']:
                new_row = [item for item in row]
                new_row.append('changed')
                diff_dress_data.append(new_row)
                continue
            if row.loc['did_you_know'] != api_data['did_you_know']:
                new_row = [item for item in row]
                new_row.append('changed')
                diff_dress_data.append(new_row)
                continue

        # check for new entries in Excel sheet that do not exist in retrieved api data
        for id in dress_ids:
            if not any(data['id'] == id for data in api_dress_data) and (sheet_dress_data['id']==id).any():
                row = sheet_dress_data.loc[sheet_dress_data['id']==id]
                new_row = [item for item in row.values[0]]
                new_row.append('new')
                diff_dress_data.append(new_row)

        column_headers = ['id', 'name', 'description', 'did_you_know', 'changed_or_new']
        generate_table(diff_dress_data, 'difference_report', column_headers, 150, 800, 'nw')

    except FileNotFoundError:
        tk.messagebox.showerror(title="Error in diffReport", message=f"File '{file_path}' not found.")
        print(f"File '{file_path}' not found.")
    except Exception as e:
        tk.messagebox.showerror(title="Error in diffReport", message=f'Error: {e}')
        print(f'Error: {e}')
    finally:
        diff_report_button.config(state="normal")

'''
Returns 3 google images url and put into excel sheet
'''
def googleImage():
    api_dress_data = sorted(apiRunner(), key=lambda x : x['id'])

    excel_file_name = "googleimages.xlsx"
    file_name = "googleimages.xlsx"
    count = 0

    # Rename file if the file name exists
    while os.path.exists(file_name):
        count += 1
        file_name = f"{os.path.splitext(excel_file_name)[0]}({count}).xlsx" # if file name exisit create new filename

    # Open the Excel file
    workbook = openpyxl.Workbook()
    sheet = workbook.active # active sheet (first sheet)

    # Create header for first row
    header_row = ["ID", "Name", "URL1", "URL2", "URL3"]
    sheet.append(header_row)

    # create progress bar
    progress_window, pb, percent_label = progress_bar('Creating Excel File')
    complete = 0

    url_list = []
    img_name_list = []

    # Loop through items in list 
    row_count = 2
    for item in api_dress_data:
        # Use google to search for image of each dress and return 3 image url
        word = f"{item['name']} india"
        url = 'https://www.google.com/search?q={0}&tbm=isch'.format(word)
        content = requests.get(url).content
        soup = BeautifulSoup(content, 'html.parser')
        images = soup.findAll('img')

        # Adds the dress number and name
        sheet.cell(row=row_count, column=1, value=item['id'])
        sheet.cell(row=row_count, column=2, value=item['name'])

        # Adds the 3 image urls to the excel sheet
        for j, image in enumerate(images[1:4]):
            img_src = image.get('src')
            try:
                response = requests.get(img_src, stream=True) # Sends a HTTP GET request to the url
                if response.status_code == 200:
                    content_type = response.headers['Content-Type'] # Gets the content type (ex: image/png)
                    img_extension = content_type.split('/')[-1].lower() # Gets the content after the "/" and put into lowercase

                    valid_extensions = {'jpg', 'jpeg', 'gif', 'png'} 
                    if img_extension in valid_extensions: # if url is a valid extension add to sheet
                        sheet.cell(row=row_count, column=j + 3, value=f"{img_src}.{img_extension}")

                        # if user chooses to download google images
                        if download_google_imgs.get() == 1: 
                            image_num = j+1
                            # creates directory to save images if one does not exist
                            if not os.path.exists('./google_images'):
                                os.makedirs('./google_images')

                            dress_name = item["name"].replace(" ", "_") # changes all space in dress name to "_"
                            new_dress_name = re.sub(r'[\\/:*?"<>|]', '', dress_name) # removes special characters
                            new_dress_name = new_dress_name.replace('\n\n', '') # removes "\n\n"

                            image_name = f'{item["id"]}_{new_dress_name}_{image_num}.{img_extension}' # sets the image file name
                            img_name_list.append(image_name)
                            url_list.append(image.get('src'))
            except Exception as e:
                print(f"Error: {e}")
        row_count+=1
        complete += 1
        pb['value'] = (complete/len(api_dress_data))*100 # calculate percentage of data retrieved
        percent_label.config(text=f'Creating Excel File...{int(pb["value"])}%') # update completion percent label

    progress_window.destroy() # close progress bar window
    workbook.save(file_name)

    if download_google_imgs.get() == 1: 
        # create progress bar
        progress_window, pb, percent_label = progress_bar('Downloading Images')
        
        threads= [] # working threads

        # spins up 10 threads at a time and calls downloadImages with url and image name
        with ThreadPoolExecutor(max_workers=10) as exec:
            for index, url in enumerate(url_list):
                threads.append(exec.submit(downloadImages, "google_images", url, img_name_list[index]))

            complete = 0 # number of threads that have finished
            for task in as_completed(threads):
                complete += 1
                pb['value'] = (complete/len(url_list))*100 # calculate percentage of images downloaded
                percent_label.config(text=f'Downloading Images...{int(pb["value"])}%') # update completion percent label
    
        progress_window.destroy() # close progress bar window

    openFile(file_name)
    google_image_search_button.config(state='normal')

'''
Performs word analysis on given dress IDs
'''
def wordAnalysis():
    word_analysis_data = [] # data from word analysis
    api_dress_data = sorted(apiRunner(), key=lambda x : x['id']) # gets dress data from API and sorts by ID

    try:
        # cycle through api_dress_data for word analysis
        for dress_data in api_dress_data:
            noun_count = 0 # number of nouns in text
            adjective_count = 0 # number of adjectives in text

            # concatenation of description and did_you_know
            text = f'{str(dress_data["description"])} {str(dress_data["did_you_know"])}'
            # TextBlob word analysis
            blob = TextBlob(text)

            # cycle through key:value pairs of TextBlob analysis to get noun and adjective count
            for k,v in blob.tags:
                if v == 'NN' or v == 'NNS' or v == 'NNP' or v == 'NNPS':
                    noun_count += 1
                elif v == 'JJ' or v == 'JJR' or v == 'JJS':
                    adjective_count += 1

            ease = textstat.flesch_reading_ease(text)
            kincaid = textstat.flesch_kincaid_grade(text)
            readability = textstat.automated_readability_index(text)

            # data to be displayed in table
            word_analysis_data.append([dress_data['id'], dress_data['name'], len(str(dress_data['description']).strip(string.punctuation).split()), 
                                       len(str(dress_data['did_you_know']).strip(string.punctuation).split()), str(noun_count), str(adjective_count),
                                       str(ease), str(kincaid), str(readability)])

        column_headers = ['id', 'name', 'description_word_count', 'did_you_know_word_count', 'total_noun_count', 'total_adjective_count', 'reading_ease', 'kincaid_grade', 'readability_index']
        generate_table(word_analysis_data, 'word_analysis_report', column_headers, 50, 200, 'center', 1)
        
    except Exception as e:
        tk.messagebox.showerror(title="Error in wordAnalysis", message=f'Error: {e}')
        print(f'Error: {e}')
    finally:
        word_analysis_button.config(state='normal')

'''
Generate Wiki Link
'''
def generateWikiLink():
    # set user agent
    USER_AGENT = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3"
    wiki_wiki = wikipediaapi.Wikipedia("en",headers={"User-Agent": USER_AGENT})

    # gather all dress data from api
    wiki_link_data = sorted(apiRunner(), key=lambda x : x['id'])
    wiki_data = []

    progress_window, pb, percent_label = progress_bar('Retrieving Wiki Data')

    for complete, item in enumerate(wiki_link_data):
        try:
            page = wiki_wiki.page(item["name"])
            if page.exists():
                item["wiki_page_link"] = page.fullurl
                wiki_data.append([item['id'], item['name'],item['wiki_page_link']])

                pb['value'] = (complete/len(wiki_link_data))*100 # calculate percentage of images downloaded
                percent_label.config(text=f'Retrieving Wiki Data...{int(pb["value"])}%') # update completion percent label
        except Exception as e:
            print(f"Error retrieving Wikipedia data for {item['name']}: {e}")

    progress_window.destroy()

    column_headers = ['id', 'name', 'wiki_page_link']
    generate_table(wiki_data, 'wiki_link_report', column_headers, 75, 800, 'nw')

    wiki_link_gen_button.config(state='normal')

'''
Generates an xcel file with the pairs found
'''
def generatePairs():
    file_path = 'APIData.xlsx'  # Change to the actual file path
    api_dress_data = sorted(apiRunner(), key=lambda x: x['id'])  # Assuming apiRunner() returns dress data
    pairs = []

    try:
        # Read dress data from Excel file
        sheet_dress_data = pd.read_excel(file_path)      
        sheet_dress_data.dropna(subset=['id'], inplace=True)  # Drop rows with missing IDs
        sheet_dress_data['description'].fillna('', inplace=True)  # Remove NA/NAN from description column
        sheet_dress_data['did_you_know'].fillna('', inplace=True)  # Remove NA/NAN from did_you_know column

	# Iterate through API dress data
        for data_that_will_be_searched in api_dress_data:
            # Get the name of the current API data ID
            name = data_that_will_be_searched['name']
            # Split the name into tokens and remove any prefixes
            tokens = [token + " " for token in name.split() if token not in ["Dr.", "Mr.", "Mrs.", "Ms."]]

            # Loop through each token
            for token in tokens:
                # Loop through provided IDs instead of the entire sheet_dress_data
                for data_that_contains_token in api_dress_data:
                    # Get the description and did you know text of the provided ID
                    description = data_that_contains_token['description']
                    did_you_know = data_that_contains_token['did_you_know']

                    # Check if the token is present in either of them, make sure we aren't looking on the same IDs
                    if data_that_will_be_searched['id'] != data_that_contains_token['id']:
                        if token in description or token in did_you_know:
                            if (data_that_contains_token['id'], data_that_will_be_searched['id']) not in [(pair[0], pair[2]) for pair in pairs]:
                                # Add the pair of IDs and names to the list
                                pairs.append([data_that_contains_token['id'], data_that_contains_token['name'], data_that_will_be_searched['id'], name])
                            # Break the inner loop as we found a pair for this token
                            break

        # Generate table and save to Excel
        column_headers = ['ID1', 'Name 1', 'ID2', 'Name 2']
        generate_table(pairs, 'generate_pairs', column_headers, 50, 200, 'center', 1)
        df_pairs = pd.DataFrame(pairs, columns=column_headers)
        df_pairs.to_excel("pairs_generated.xlsx", index=False)
        print("Excel file 'pairs_generated.xlsx' created.")

    except FileNotFoundError:
        print(f"File '{file_path}' not found.")
    except Exception as e:
        print(f'Error: {e}')

'''
Function to fetch the english text from dress data
'''
def fetch_english_text(dress_data):
    english_texts = {}  
    for dress in dress_data:
        dress_id = dress.get('id')
        english_description = {
            'name': dress.get('name', ''),
            'description': dress.get('description', ''),
            'did_you_know': dress.get('did_you_know', '')
        }
        english_texts[dress_id] = english_description
    return english_texts
'''
translation function that stores the translated telugu texts into the same "keys" as the english for easier formatting
'''
def translate_text_to_telugu(english_texts):
    translator = googletrans.Translator()
    telugu_texts = {}

    for id, texts in english_texts.items():
        telugu_texts[id] = {}  # Initialize a dictionary for this ID
        for key, text in texts.items():
            try:
                translated = translator.translate(text, dest='te')  # 'te' for Telugu
                telugu_texts[id][key] = translated.text
            except Exception as e:
                print(f"Error translating {key} for ID {id}: {e}")
                telugu_texts[id][key] = text  # Use the original text if translation fails

    return telugu_texts

def translate_text_to_first_person(english_texts):
    first_person_texts = {} # Dictionary to hold the translated text
    messages = [{"role": "system", "content": 
                 "You translate text from third person to first person. Don't reply with anything. Just do the work."}]


    for id, texts in english_texts.items():
        first_person_texts[id] = {}
        for key, text in texts.items():

            # Perform conversion from third person to first person using ChatGPT
            messages.append(
                    {"role": "user", "content": f"Convert the following text from the third person to the first person?\n\n{text}\n\n"}
                    )
            chat = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=messages
            )
            reply = chat.choices[0].message.content
            first_person_texts[id][key] = reply
            print(f"This is the Reply: {reply}\n\n")
            print(first_person_texts)
            #first_person_texts[id] = first_person_text.choices[0].text.strip()
    return first_person_texts

"""
    Create an HTML package with English and Telugu texts.
    """
def create_html_package(english_texts, telugu_texts):
    page = 1
    html_content = """
    <html>
    <head>
    <style>
        .name { font-weight: bold; text-align: center; }
        .did_you_know { margin-top: 20px; font-size: 18px; }
        .page {
            page-break-after: always;
        }
        @media screen {
            .page {
                border-bottom: 1px solid #ccc;
                padding-bottom: 20px;
                margin-bottom: 20px;
            }
        }
    </style>
    </head>
    <body>
    """
    for id, english_text in english_texts.items():
        telugu_text = telugu_texts.get(id, {'name': '', 'description': '', 'did_you_know': ''})
        html_content += f"<div class='page'><h2>Page No: {page} ABCDid: {id} (English)</h2>"
        html_content += f"<div class='name'>{english_text['name']}</div>"
        html_content += f"<p class='description'>{english_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{english_text['did_you_know']}</p></div><hr>"
        
        html_content += f"<div class='page'><h2>Page No: {page} ABCDid: {id} (Telugu)</h2>"
        html_content += f"<div class='name'>{telugu_text['name']}</div>"
        html_content += f"<p class='description'>{telugu_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{telugu_text['did_you_know']}</p></div><hr>"
        page += 1
    html_content += "</body></html>"
    return html_content

def create_html_package_gpt(english_texts, first_person_texts):
    page = 1
    html_content = """
    <html>
    <head>
    <style>
        .name { font-weight: bold; text-align: center; }
        .did_you_know { margin-top: 20px; font-size: 18px; }
        .page {
            page-break-after: always;
        }
        @media screen {
            .page {
                border-bottom: 1px solid #ccc;
                padding-bottom: 20px;
                margin-bottom: 20px;
            }
        }
    </style>
    </head>
    <body>
    """
    for id, english_text in english_texts.items():
        first_person_text = first_person_texts.get(id, {'name': '', 'description': '', 'did_you_know': ''})
        html_content += f"<div class='page'><h2>Page No: {page} ABCD-id: {id} (English)</h2>"
        html_content += f"<div class='name'>{english_text['name']}</div>"
        html_content += f"<div class='description>'>Description</div>"
        html_content += f"<p>{english_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{english_text['did_you_know']}</p></div><hr>"
        
        html_content += f"<div class='page'><h2>Page No: {page} ABCD-id: {id} (ChatGPT)</h2>"
        html_content += f"<div class='name'>{english_text['name']}</div>"
        html_content += f"<div class='description>'>Description</div>"
        html_content += f"<p>{first_person_text['description']}</p>"
        html_content += "<div class='did_you_know'>Did you know?</div>"
        html_content += f"<p>{first_person_text['did_you_know']}</p></div><hr>"
        page += 1
    html_content += "</body></html>"
    return html_content

"""
    Save the HTML content to a file, appending an incrementing number to the filename
    to avoid overwrites.
    """
def save_html_to_file(english_texts, telugu_texts, base_filename="translation_package"):
    html_filename = f"{base_filename}.html"
    txt_filename = f"{base_filename}.txt"
    counter = 1
    # Check if the file exists and update the filename until it's unique
    while os.path.exists(html_filename) or os.path.exists(txt_filename):
        html_filename = f"{base_filename}_{counter}.html"
        txt_filename = f"{base_filename}_{counter}.txt"
        counter += 1
    
    # Saving HTML content
    with open(html_filename, 'w', encoding='utf-8') as file:
        file.write(create_html_package(english_texts, telugu_texts))  # Assuming create_html_package is defined as before
    
    # Creating a plain text version of the content
    text_content = ""
    for id, english_text in english_texts.items():
        telugu_text = telugu_texts.get(id, '')
        text_content += f"English Text for ID {id}\n{english_text}\n\n"
        text_content += "------------------------------------------------\n\n"
        text_content += f"Telugu Text for ID {id}\n{telugu_text}\n\n"
        text_content += "================================================\n\n"
    
    # Saving text content
    with open(txt_filename, 'w', encoding='utf-8') as file:
        file.write(text_content)

    return html_filename, txt_filename

def save_html_to_file_gpt(english_texts, first_person_texts, base_filename="gpt_adjusted_package"):
    html_filename = f"{base_filename}.html"
    txt_filename = f"{base_filename}.txt"
    counter = 1
    # Check if the file exists and update the filename until it's unique
    while os.path.exists(html_filename) or os.path.exists(txt_filename):
        html_filename = f"{base_filename}_{counter}.html"
        txt_filename = f"{base_filename}_{counter}.txt"
        counter += 1
    
    # Saving HTML content
    with open(html_filename, 'w', encoding='utf-8') as file:
        file.write(create_html_package_gpt(english_texts, first_person_texts))  # Assuming create_html_package is defined as before
    
    # Creating a plain text version of the content
    text_content = ""
    for id, english_text in english_texts.items():
        telugu_text = first_person_texts.get(id, '')
        text_content += f"English Text for ID {id}\n{english_text}\n\n"
        text_content += "------------------------------------------------\n\n"
        text_content += f"ChatGPT Text for ID {id}\n{telugu_text}\n\n"
        text_content += "================================================\n\n"
    
    # Saving text content
    with open(txt_filename, 'w', encoding='utf-8') as file:
        file.write(text_content)

    return html_filename, txt_filename

def generate_translation_package():
    # Get ID numbers from text area
    try:

        dress_data = apiRunner()
        
        # Fetch English text
        english_texts = fetch_english_text(dress_data)
        

        # Translate to Telugu
        telugu_texts = translate_text_to_telugu(english_texts)
        
        # Create HTML package
        html_content = create_html_package(english_texts, telugu_texts)
    
        
        # Save HTML and TXT files
        html_filename, txt_filename = save_html_to_file(english_texts, telugu_texts)
    
        
        # Open the HTML file in a web browser
        webbrowser.open(f'file://{os.path.realpath(html_filename)}')
    
    except FileNotFoundError:
        print(f"File '{e}' not found.")
    except Exception as e:
        print(f'Error: {e}')
    finally:
        translation_package_generate_button.config(state='normal')
        
        # Optionally, show a message that the file has been saved
        print("Translation package has been generated and saved.")



def generate_first_person_package():
    try:

        dress_data = apiRunner() 
        english_texts = fetch_english_text(dress_data)
        first_person_texts = translate_text_to_first_person(english_texts)
        html_content = create_html_package_gpt(english_texts, first_person_texts)
        html_filename, txt_filename = save_html_to_file_gpt(english_texts, first_person_texts)
        webbrowser.open(f'file://{os.path.realpath(html_filename)}')
    except FileNotFoundError:
        print(f"File '{e}' not found")
    except Exception as e:
        print(f'Erro: {e}')
    finally:
        first_person_generate_button.config(state='normal')
    print("Translation package has been generated and saved.")

def generate_first_person_package():
    dress_data = apiRunner() 
    english_texts = fetch_english_text(dress_data)
    first_person_texts = translate_text_to_first_person(english_texts)
    html_content = create_html_package(english_texts, first_person_texts)
    html_filename, txt_filename = save_html_to_file(html_content)
    webbrowser.open(f'file://{os.path.realpath(html_filename)}')
    print("Translation package has been generated and saved.")

def wordSearchOpenAi(english_texts):
    words_for_puzzles = {}
    word_count = int(word_count_var.get())  # Get the current preferred word count

    messages = [
        {"role": "system", "content": f"You will extract {word_count} meaningful words significant to the character's text. Do not reply with anything else. Just list the words."}
    ]

    #sorted_texts = sort_english_texts(english_texts)  # Sort texts based on preference

    for id, text in english_texts.items():
        messages.append(
            {"role": "user", "content": f"Can you extract {word_count} words from this character's text that are meaningful and significant to the character, with one being their name? Please only output the {word_count} words.\n\n{text}\n\n"}
        )
        
        chat_response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=messages
        )
       
        reply = chat_response.choices[0].message.content.strip()
        
        words_for_puzzles[id] = reply
        print(f"This is the Reply for ID {id}: {reply}\n\n")
    
    return words_for_puzzles


#function that splits the words and ids, preparing them for puzzle creation
def wordsearchCreator(words_for_puzzles):
    puzzles = {}
    answer_keys = {}  
    word_lists = {}
    for id, words_string in words_for_puzzles.items():
        word_list = words_string.replace(',', '').replace("'", "").upper().split()
        word_lists[id] = word_list
        grid_size = int(puz_width_var.get())  # Ensure this is correctly fetched
        grid = [['-' for _ in range(grid_size)] for _ in range(grid_size)]
        answer_positions = {}

        for word in word_list:
            placed, positions = placeWord(grid, word)  # This function needs to be adapted to return start and end points too
            if placed:
                answer_positions[word] = {'start': positions[0], 'end': positions[-1]}

        fillEmptySpots(grid)
        puzzles[id] = grid
        answer_keys[id] = answer_positions
    
    return puzzles, answer_keys, word_lists


#function to randomly place words in the grid
def placeWord(grid, word):
    max_attempts = 100    
    attempts = 0
    placed = False
    positions = []  

    while not placed and attempts < max_attempts:
        wordPlacement = random.randint(0, 3)
        attempts += 1

        if wordPlacement == 0:  # Horizontal
            row = random.randint(0, len(grid) - 1)
            col = random.randint(0, len(grid) - len(word))
            space_available = all(grid[row][col + i] == '-' or grid[row][col + i] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row][col + i] = word[i]
                    positions.append((row, col + i))  
                placed = True

        elif wordPlacement == 1:  # Vertical
            row = random.randint(0, len(grid) - len(word))
            col = random.randint(0, len(grid) - 1)
            space_available = all(grid[row + i][col] == '-' or grid[row + i][col] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row + i][col] = word[i]
                    positions.append((row + i, col))  
                placed = True

        elif wordPlacement == 2:  # Diagonal left to right
            row = random.randint(0, len(grid) - len(word))
            col = random.randint(0, len(grid) - len(word))
            space_available = all(grid[row + i][col + i] == '-' or grid[row + i][col + i] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row + i][col + i] = word[i]
                    positions.append((row + i, col + i)) 
                placed = True

        elif wordPlacement == 3:  # Diagonal right to left
            row = random.randint(0, len(grid) - len(word))
            col = random.randint(len(word) - 1, len(grid) - 1)
            space_available = all(grid[row + i][col - i] == '-' or grid[row + i][col - i] == word[i] for i in range(len(word)))
            if space_available:
                for i in range(len(word)):
                    grid[row + i][col - i] = word[i]
                    positions.append((row + i, col - i))  
                placed = True

    return placed, positions

#function to fill the empty slots after words are placed                 
def fillEmptySpots(grid):
    for row in range(len(grid)):
        for col in range(len(grid[0])):  
            if grid[row][col] == '-':  
                grid[row][col] = random.choice(string.ascii_uppercase)
            
def createWordsearchWordsHtml(puzzles, answer_keys, word_lists):
    page = 1
    html_content = """
    <html>
    <head>
    <title>Word Search Puzzles</title>
    <style>
        body {
            font-family: Arial, sans-serif;
        }
        .page {
            page-break-after: always;
            margin-bottom: 20px;
        }
        table {
            border-collapse: collapse;
            margin: 20px 0;
            position: relative;
        }
        td {
            border: 1px solid #666;
            width: 20px;
            height: 20px;
            text-align: center;
            vertical-align: middle;
        }
        .highlighted {
            font-weight: bold;  /* Make the text bold */
            color: black;  /* Ensure the text is black */
        }
        .answer-line {
            position: absolute;
            stroke: red;
            stroke-width: 2;
            marker-end: url(#arrowhead);
        }
        .svg-container {
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            pointer-events: none;  /* Allows clicks to pass through to the table */
        }
    </style>
    </head>
    <body>
    <svg style="display:none;">
        <defs>
            <marker id="arrowhead" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto" markerUnits="strokeWidth">
                <path d="M0,0 L0,6 L9,3 z" fill="red" />
            </marker>
        </defs>
    </svg>
    """

    # Generate regular puzzles with word lists
    for id, grid in puzzles.items():
        html_content += f"<div class='page'><h2>Page No: {page} - Puzzle ID: {id}</h2><table>"
        for row in grid:
            html_content += "<tr>"
            for cell in row:
                html_content += f"<td>{cell}</td>"
            html_content += "</tr>"
        html_content += "</table><div><strong>Words:</strong><ul>"
        for word in word_lists[id]:
            html_content += f"<li>{word}</li>"
        html_content += "</ul></div></div>"
        page += 1

    # Generate answer key puzzles and draw SVG lines directly for each word
    for id, positions in answer_keys.items():
        grid = puzzles[id]
        html_content += f"<div class='page'><h2>Answer Key Page No: {page} - Puzzle ID: {id}</h2><div style='position: relative;'>"
        html_content += "<table>"
        for row_idx, row in enumerate(grid):
            html_content += "<tr>"
            for col_idx, cell in enumerate(row):
                if (row_idx, col_idx) in positions:
                    html_content += f"<td class='highlighted'>{cell}</td>"
                else:
                    html_content += f"<td>{cell}</td>"
            html_content += "</tr>"
        html_content += "</table>"

        # SVG overlay for drawing lines
        html_content += "<div class='svg-container'><svg style='width: 100%; height: 100%;'>"
        for word, pos in positions.items():
            start = pos['start']
            end = pos['end']
            start_x = start[1] * 20 + 10
            start_y = start[0] * 20 + 10
            end_x = end[1] * 20 + 10
            end_y = end[0] * 20 + 10
            html_content += f"<line x1='{start_x}' y1='{start_y}' x2='{end_x}' y2='{end_y}' class='answer-line'></line>"
        html_content += "</svg></div></div>"
        page += 1

    html_content += "</body></html>"
    return html_content







def save_and_display_html(html_content, base_filename="puzzles_package"):
    html_filename = f"{base_filename}.html"
    counter = 1
    
    # Increment filename if exists to avoid overwriting
    while os.path.exists(html_filename):
        html_filename = f"{base_filename}_{counter}.html"
        counter += 1

    # Save HTML to file
    with open(html_filename, 'w') as file:
        file.write(html_content)
    print(f"HTML content has been saved to {html_filename}.")
    
    # Format the file path for browser compatibility and open it
    try:
        file_url = f"file://{os.path.abspath(html_filename)}"
        webbrowser.open(file_url, new=2)
        print("HTML file has been opened in your web browser.")
    except Exception as e:
        print(f"Failed to open the HTML file in a web browser. Error: {e}")
        
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR


def add_puzzle_table(slide, grid, word_list, title_text, answer_positions=None):
    title = slide.shapes.title
    title.text = title_text

    rows, cols = len(grid), len(grid[0])
    grid_origin_x = Inches(1)  
    grid_origin_y = Inches(1.5)  
    max_width = Inches(6)  
    max_height = Inches(4.5)  # Total available height for the grid
    cell_width = max_width / cols
    cell_height = max_height / rows

    # Add the grid table
    table = slide.shapes.add_table(rows, cols, grid_origin_x, grid_origin_y, round(cell_width * cols), round(cell_height * rows)).table
    table.first_row = False

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = grid[r][c]
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add lines for correct answers
    if answer_positions:
        for word, pos in answer_positions.items():
            start_cell = pos['start']
            end_cell = pos['end']
            start_x = grid_origin_x + start_cell[1] * cell_width + cell_width / 2
            start_y = grid_origin_y + start_cell[0] * cell_height + cell_height / 2
            end_x = grid_origin_x + end_cell[1] * cell_width + cell_width / 2
            end_y = grid_origin_y + end_cell[0] * cell_height + cell_height / 2

            line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, start_x, start_y, end_x, end_y)
            line.line.width = Pt(2)
            line.line.color.rgb = RGBColor(255, 0, 0)  # Set line color to red

    # Add word list to the side
    textbox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(2), Inches(4))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "Words to find:\n" + "\n".join(word_list)
    p.font.bold = True
    p.font.size = Pt(14)



def make_powerpoint(puzzles, answer_keys, word_lists):
    prs = Presentation()
    for id, grid in puzzles.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank layout
        add_puzzle_table(slide, grid, word_lists[id], f"Puzzle ID: {id}")

    # Add a slide for each answer key
    for id, positions in answer_keys.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank layout
        add_puzzle_table(slide, puzzles[id], word_lists[id], f"Answer Key ID: {id}", answer_positions=positions)

    return prs




def save_powerpoint(prs, base_filename="puzzles_package"):
    filename = f"{base_filename}.pptx"
    counter = 1
    
    # Increment filename if exists to avoid overwriting
    while os.path.exists(filename):
        filename = f"{base_filename}_{counter}.pptx"
        counter += 1

    # Save PowerPoint to file
    prs.save(filename)
    print(f"PowerPoint content has been saved to {filename}.")
    
def generate_word_search_package():
    
    dress_data = apiRunner()

    
    english_texts = fetch_english_text(dress_data)

   
    puzzleWords = wordSearchOpenAi(english_texts)

    
    puzzles, answer_keys, word_lists = wordsearchCreator(puzzleWords)

    
    html_puzzles = createWordsearchWordsHtml(puzzles, answer_keys, word_lists)

   
    save_and_display_html(html_puzzles)
    print("HTML word puzzles have been generated and saved.")

   
    prs = make_powerpoint(puzzles, answer_keys, word_lists)

    
    save_powerpoint(prs, base_filename="word_puzzles_package")
    print("PowerPoint word puzzles have been generated and saved.")



'''
Spins up new thread to run generateUpdate function
'''
def startGenerateBookThread():
    book_gen_generate_button.config(state="disabled")
    generate_thread = threading.Thread(target=generateBook)
    generate_thread.start()

'''
Spins up new thread to run diffReport function
'''
def startDiffReportThread():
    diff_report_button.config(state='disabled')
    diff_report_thread = threading.Thread(target=diffReport)
    diff_report_thread.start()

'''
Spins up new thread to run wordAnalysis function
'''
def startWordAnalysisThread():
    word_analysis_button.config(state='disabled')
    word_analysis_thread = threading.Thread(target=wordAnalysis)
    word_analysis_thread.start()

'''
Spins up new thread to run googleImage function
'''
def startGoogleImageThread():
    google_image_search_button.config(state='disabled')
    google_image_search_thread = threading.Thread(target=googleImage)
    google_image_search_thread.start()

'''
Spins up new thread to run generateWikiLink function
'''
def startGenerateWikiLinkThread():
    wiki_link_gen_button.config(state='disabled')
    wiki_link_thread = threading.Thread(target=generateWikiLink)
    wiki_link_thread.start()

'''
Spins up new thread to run generatePairs function
'''
def startGeneratePairsThread():
    who_are_my_pairs_gen_button.config(state='disabled')
    who_are_my_pairs_thread = threading.Thread(target=generatePairs)
    who_are_my_pairs_thread.start()

'''
Spins up new thread to run translatepackage function
'''
def startTranslationPackageThread():
    translation_package_generate_button.config(state="disabled")
    translate_package_thread = threading.Thread(target=generate_translation_package)
    translate_package_thread.start()

'''
Spins up new thread to run translate_to_first person function
'''
def startFirstPersonThread():
    first_person_generate_button.config(state="disabled")
    first_person_thread = threading.Thread(target=generate_first_person_package)
    first_person_thread.start()

'''
Spins up new thread to run word search function
'''
def startWordPuzzleThread():
    word_puzzle_generate_button.config(state="disabled")
    word_search_thread = threading.Thread(target=generate_word_search_package)
    word_search_thread.start()

'''
Launch help site when user clicks Help button
'''
def launchHelpSite():
    # create help site
    with open('help.html', 'w') as file:
        file.write('<!DOCTYPE html>\n<html>\n<head>\n\t<meta charset="utf8">\n\t<title>abcd Help</title>\n</head>\n<body>\n\t\t<h1 style="text-align: center;">Welcome to the help site</h1>\n</body>\n</html>\n')
    
    # open help site
    webbrowser.open('help.html')

'''
Raise selected frame to the top
'''
def raiseFrame(frame):
    if frame == 'main_frame':
        main_frame.tkraise()
        root.title("Project ABCD Admin Panel")
    elif frame == 'book_gen_frame':
        book_gen_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Book Generation")
    elif frame == 'diff_report_frame':
        diff_report_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Difference Report")
    elif frame == 'word_analysis_frame':
        word_analysis_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Analysis")
    elif frame == 'google_image_frame':
        google_image_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Google Image")
    elif frame == 'wiki_link_frame':
        wiki_link_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Wiki Link")
    elif frame == 'who_are_my_pairs_frame':
        who_are_my_pairs_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Who Are My Pairs")
    elif frame == 'translation_package_frame':  
        translation_package_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Translation Package")
    elif frame == 'first_person_frame':
        first_person_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
    elif frame == 'word_puzzle_frame':
        word_puzzle_frame.tkraise()
        text_field_label.tkraise()
        text_field.tkraise()
        root.title("Project ABCD Word Search Puzzles")
        


#--------------------------------Main Frame-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------------------------------------------------
# configure grid to fill extra space and center
tk.Grid.rowconfigure(root, 0, weight=1)
tk.Grid.columnconfigure(root, 0, weight=1)

# main frame
main_frame = tk.Frame(root, width=1000, height=600)
main_frame.pack_propagate(False)
main_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Main Title-----------------------------------------------------------------------------------------------
# Create a label widget
title_label = tk.Label(main_frame, text="Project ABCD\nMain Menu",  font=('Arial', 20))
title_label.pack(pady=100)

#--------------------------------Main Buttons-----------------------------------------------------------------------------------------------
# Create buttons widget
## Button settings
main_button_frame = tk.Frame(main_frame)
main_button_frame.place(relx=.5, rely=.5, anchor='center')
button_width = 20
button_height = 3
button_bgd_color = "#007FFF"
button_font_color = "#ffffff"

## Generate Book: Gets selected dress from API and import into ppt
generate_book_button = tk.Button(main_button_frame, text="Generate Book", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('book_gen_frame'))
generate_book_button.pack(side="left", padx=50)

## Diff Report: Create a SQL file of dresses that got changed from excel sheet byt comparing to API
diff_report_button = tk.Button(main_button_frame, text="Difference Report", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('diff_report_frame'))
diff_report_button.pack(side="left", padx=50, anchor='center')

## Generate Book: Get selected dress that user input & put into a table (ID, Name, Description Count, DYK Count, Total Nouns Count, Total Adjectives Count)
word_analysis_report_button = tk.Button(main_button_frame, text="Word Analysis Report", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('word_analysis_frame'))
word_analysis_report_button.pack(side="left", padx=50)

main_button_frame2 = tk.Frame(main_frame)
main_button_frame2.place(relx=.5, rely=.7, anchor='center')

## Google Images: Create an Excel file with 3 image links to the selected dresses
google_image_button = tk.Button(main_button_frame2, text="Google Image", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('google_image_frame'))
google_image_button.pack(side="left", padx=50)

## Wiki Link: [FILL IN THE ACTION HERE]
wiki_link_button = tk.Button(main_button_frame2, text="Wiki Link", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('wiki_link_frame'))
wiki_link_button.pack(side="left", padx=50)

## My Pairs: Shows pairs when they are searched
who_are_my_pairs_button = tk.Button(main_button_frame2, text="Who Are My Pairs?", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('who_are_my_pairs_frame'))
who_are_my_pairs_button.pack(side="left", padx=50)

main_button_frame3 = tk.Frame(main_frame)
main_button_frame3.place(relx=.5, rely=.9, anchor='center')

## Generate Book: fetches English text from Api, uses google translate to generate "telugu" text, then creates HTML package with english text on page, and "telugu" text on another.
translation_package_button = tk.Button(main_button_frame3, text="Translation Package", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('translation_package_frame'))
translation_package_button.pack(side="left", padx=50)

## First Person: fetches text from api, uses ChatGPT to reword the description and did you know text to first person
first_person_button = tk.Button(main_button_frame3, text="First Person Conversion", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('first_person_frame'))
first_person_button.pack(side="left", padx=50)

## Word Puzzle: generates and creates crossword puzzles based of words in character descriptions
word_puzzle_button = tk.Button(main_button_frame3, text="Word Puzzle Creator", font=LABEL_FONT, width=button_width, height=button_height, bg=button_bgd_color, fg=button_font_color, command=lambda: raiseFrame('word_puzzle_frame'))
word_puzzle_button.pack(side="left", padx=50)
#--------------------------------Book Gen Frame---------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------------------------------------------------
# book gen frame
book_gen_frame = tk.Frame(root, width=1000, height=600)
book_gen_frame.pack_propagate(False)
book_gen_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Text Field-----------------------------------------------------------------------------------------------
# text field label
text_field_label = tk.Label(root, text="Dress Numbers:", font=LABEL_FONT)
text_field_label.place(x=25, y=72.5)

# text field
text_field = tk.Text(root)
text_field.place(x=175, y=10, relwidth=.8, height=135)

# text field initialization
try:
    with open('slide_numbers.txt', 'r') as file:
        slide_number_content = file.readline().strip()
        text_field.insert("1.0", slide_number_content)
except FileNotFoundError:
    print(FileNotFoundError)

#--------------------------------Layout Radio Buttons-------------------------------------------------------------------------------------
# layout variable
layout = tk.IntVar()
layout.set(4)

# layout frame
layout_frame = tk.Frame(book_gen_frame)
# layout radio buttons
layout_radio4 = tk.Radiobutton(layout_frame, text="Picture on Left - Text on right - Single Page - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=4)
layout_radio1 = tk.Radiobutton(layout_frame, text="Picture on Left Page - Text on Right Page - Two Page Mode - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=1)
layout_radio2 = tk.Radiobutton(layout_frame, text="Picture on Right - Text on Left - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=2)
layout_radio3 = tk.Radiobutton(layout_frame, text="Picture on Left - Text on Right - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=3)
# pack radio buttons into layout frame
layout_radio4.pack(anchor="nw")
layout_radio1.pack(anchor="nw")
layout_radio2.pack(anchor="nw")
layout_radio3.pack(anchor="nw")
# place layout frame on main frame
layout_frame.place(x=175, y=150, width=800)
# layout radio buttons label
layout_label = tk.Label(book_gen_frame, text="Layout:", font=LABEL_FONT)
layout_label.place(x=25, y=170)

# separator line
separator1 = ttk.Separator(book_gen_frame)
separator1.place(x=175, y=150, relwidth=.8)

#--------------------------------Sort Radio Buttons---------------------------------------------------------------------------------------
# sort variable
sort_order = tk.IntVar()
sort_order.set(1)

# sort frame
sort_frame = tk.Frame(book_gen_frame)
# sort radio buttons
sort_radio1 = tk.Radiobutton(sort_frame, text="By Name", font=MAIN_FONT, variable=sort_order, value=1)
sort_radio2 = tk.Radiobutton(sort_frame, text="By ID", font=MAIN_FONT, variable=sort_order, value=2)
sort_radio3 = tk.Radiobutton(sort_frame, text="By Input Order", font=MAIN_FONT, variable=sort_order, value=3)
# pack radio buttons into sort frame
sort_radio1.pack(side="left")
sort_radio2.pack(side="left")
sort_radio3.pack(side="left")
# place sort frame on main frame
sort_frame.place(x=175, y=265, width=800)

# sort radio buttons label
sort_label = tk.Label(book_gen_frame, text="Sort Order:", font=LABEL_FONT)
sort_label.place(x=25, y=265)


# separator line
separator2 = ttk.Separator(book_gen_frame)
separator2.place(x=175, y=265, relwidth=.8)


#--------------------------------Preferences----------------------------------------------------------------------------------------------
# Initialize variables for storing values
text_size_var = tk.StringVar()
title_size_var = tk.StringVar()
subtitle_size_var = tk.StringVar()
text_font_var = tk.StringVar()
title_font_var = tk.StringVar()
subtitle_font_var = tk.StringVar()
pic_width_var = tk.StringVar()
pic_height_var = tk.StringVar()

# Set initial values for the Entry fields
text_size_var.set(preferences["TEXT_SIZE"])
title_size_var.set(preferences["TITLE_SIZE"])
subtitle_size_var.set(preferences["SUBTITLE_SIZE"])
text_font_var.set(preferences["TEXT_FONT"])
title_font_var.set(preferences["TITLE_FONT"])
subtitle_font_var.set(preferences["SUBTITLE_FONT"])
pic_width_var.set(preferences["PIC_WIDTH"])
pic_height_var.set(preferences["PIC_HEIGHT"])

# preferences frame
preferences_frame = tk.Frame(book_gen_frame)
# preferences labels and entry fields
text_size_label = tk.Label(preferences_frame, text="Text Size:", font=MAIN_FONT)
text_size = tk.Entry(preferences_frame, width=3, textvariable=text_size_var, state="disabled", font=MAIN_FONT)

title_size_label = tk.Label(preferences_frame, text="Title Size:", font=MAIN_FONT)
title_size = tk.Entry(preferences_frame, width=3, textvariable=title_size_var, state="disabled",  font=MAIN_FONT)

subtitle_size_label = tk.Label(preferences_frame, text="Subtitle Size:", font=MAIN_FONT)
subtitle_size = tk.Entry(preferences_frame, width=3, textvariable=subtitle_size_var, state="disabled",  font=MAIN_FONT)

text_font_label = tk.Label(preferences_frame, text="Text Font:", font=MAIN_FONT)
text_font = tk.Entry(preferences_frame, width=25, textvariable=text_font_var, state="disabled",  font=MAIN_FONT)

title_font_label = tk.Label(preferences_frame, text="Title Font:", font=MAIN_FONT)
title_font = tk.Entry(preferences_frame, width=25, textvariable=title_font_var, state="disabled",  font=MAIN_FONT)

subtitle_font_label = tk.Label(preferences_frame, text="Subitle Font:", font=MAIN_FONT)
subtitle_font = tk.Entry(preferences_frame, width=25, textvariable=subtitle_font_var, state="disabled",  font=MAIN_FONT)

pic_width_label = tk.Label(preferences_frame, text="Pic Width:", font=MAIN_FONT)
pic_width = tk.Entry(preferences_frame, width=6, textvariable=pic_width_var, state="disabled",  font=MAIN_FONT)

pic_height_label = tk.Label(preferences_frame, text="Pic Height:", font=MAIN_FONT)
pic_height = tk.Entry(preferences_frame, width=6, textvariable=pic_height_var, state="disabled",  font=MAIN_FONT)

# grid preference labels and entry fields into preferences frame
# column 1 + 2
text_size_label.grid(row=1, column=1, pady=10)
text_size.grid(row=1, column=2)

title_size_label.grid(row=2, column=1, pady=10)
title_size.grid(row=2, column=2)

subtitle_size_label.grid(row=3, column=1, pady=10, padx=15)
subtitle_size.grid(row=3, column=2)

# column 3 + 4
text_font_label.grid(row=1, column=3)
text_font.grid(row=1, column=4)

title_font_label.grid(row=2, column=3)
title_font.grid(row=2, column=4)

subtitle_font_label.grid(row=3, column=3, padx=15)
subtitle_font.grid(row=3, column=4)

# column 5 + 6
pic_width_label.grid(row=1, column=5)
pic_width.grid(row=1, column=6)

pic_height_label.grid(row=2, column=5, padx=15)
pic_height.grid(row=2, column=6)

# place preferences frame on main frame
preferences_frame.place(x=175, y=295, width=800)

# preferences label
preferences_label = tk.Label(book_gen_frame, text="Preferences:", font=LABEL_FONT)
preferences_label.place(x=25, y=350)

# separator line
separator3 = ttk.Separator(book_gen_frame)
separator3.place(x=175, y=295, relwidth=.8)

#--------------------------------Numbering Radio Buttons--------------------------------------------------------------------------------------
# numbering variable
numbering = tk.IntVar()
numbering.set(1)

# numbering frame
numbering_frame = tk.Frame(book_gen_frame)
# numbering radio buttons
numbering_radio1 = tk.Radiobutton(numbering_frame, text="Show both Page Number and Dress ID", font=MAIN_FONT, variable=numbering, value=1)
numbering_radio2 = tk.Radiobutton(numbering_frame, text="Show Page Number", font=MAIN_FONT, variable=numbering, value=2)
numbering_radio3 = tk.Radiobutton(numbering_frame, text="Show Dress ID", font=MAIN_FONT, variable=numbering, value=3)
# pack numbering buttons into sort frame
numbering_radio1.pack(side="left")
numbering_radio2.pack(side="left")
numbering_radio3.pack(side="left")
# place numbering frame on main frame
numbering_frame.place(x=175, y=445, width=800)

# numbering radio buttons label
numbering_label = tk.Label(book_gen_frame, text="Numbering:", font=LABEL_FONT)
numbering_label.place(x=25, y=445)

# separator line
separator4 = ttk.Separator(book_gen_frame)
separator4.place(x=175, y=445, relwidth=.8)

#--------------------------------Translate and Image Check Buttons--------------------------------------------------------------------------------------
# translate check variable
translate = tk.IntVar()
translate.set(0)
# language options variable
language = tk.StringVar()
language.set(LANGUAGES[0])
# Download image variable
download_imgs = tk.IntVar()
download_imgs.set(0)
# Generate book from local Excel sheet
gen_local = tk.IntVar()
gen_local.set(0)

# translate frame
check_button_frame = tk.Frame(book_gen_frame)
# translate check button
translate_checkbutton = tk.Checkbutton(check_button_frame, text="Translate to:", font=MAIN_FONT, variable=translate, onvalue=1, offvalue=0)
# language options
language_options = tk.OptionMenu(check_button_frame, language, *LANGUAGES)
# download images
download_images = tk.Checkbutton(check_button_frame, text="Download Images", font=MAIN_FONT, variable=download_imgs, onvalue=1, offvalue=0, command=lambda: gen_local.set(0))
# generate book from local Excel sheet
generate_from_local = tk.Checkbutton(check_button_frame, text="Generate from Local", font=MAIN_FONT, variable=gen_local, onvalue=1, offvalue=0, command=lambda: download_imgs.set(0))

# pack translate options into translate frame
translate_checkbutton.pack(side="left")
language_options.pack(side="left")
download_images.pack(side="left")
generate_from_local.pack(side="left")
# place translate frame on main frame
check_button_frame.place(x=175, y=495)

# separator line
separator5 = ttk.Separator(book_gen_frame)
separator5.place(x=175, y=495, relwidth=.8)

#--------------------------------Book Gen Buttons--------------------------------------------------------------------------------------
# button frame
book_gen_button_frame = tk.Frame(book_gen_frame)
# generate button
book_gen_generate_button = tk.Button(book_gen_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGenerateBookThread)
# help button
book_gen_help_button = tk.Button(book_gen_button_frame, text="Help", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=launchHelpSite)
# upload button
book_gen_back_button = tk.Button(book_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
book_gen_generate_button.pack(side="left", padx=35)
book_gen_help_button.pack(side="left")
book_gen_back_button.pack(side="left", padx=30)

# place button frame on main frame
book_gen_button_frame.pack(side="bottom", pady=10)


#--------------------------------Diff Report Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
diff_report_frame = tk.Frame(root, width=1000, height=600)
diff_report_frame.pack_propagate(False)
diff_report_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Diff Report Buttons-----------------------------------------------------------------------------------------------
# button frame
diff_report_button_frame = tk.Frame(diff_report_frame)
# difference report button
diff_report_button = tk.Button(diff_report_button_frame, text="Diff Report", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startDiffReportThread)
# back button
diff_back_button = tk.Button(diff_report_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
diff_report_button.pack(side="left", padx=35)
diff_back_button.pack(side="left", padx=30)

# place button frame on diff report frame
diff_report_button_frame.pack(side="bottom", pady=10)


#--------------------------------Word Analysis Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
word_analysis_frame = tk.Frame(root, width=1000, height=600)
word_analysis_frame.pack_propagate(False)
word_analysis_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Word Analysis Buttons-----------------------------------------------------------------------------------------------
# button frame
word_analysis_button_frame = tk.Frame(word_analysis_frame)
# word analysis button
word_analysis_button = tk.Button(word_analysis_button_frame, text="Word Analysis", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startWordAnalysisThread)
# back button
word_analysis_back_button = tk.Button(word_analysis_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
word_analysis_button.pack(side="left", padx=35)
word_analysis_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
word_analysis_button_frame.pack(side="bottom", pady=10)

#--------------------------------Google Image Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
google_image_frame = tk.Frame(root, width=1000, height=600)
google_image_frame.pack_propagate(False)
google_image_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Google Image Buttons-----------------------------------------------------------------------------------------------
# Download google image variable
download_google_imgs = tk.IntVar()
download_google_imgs.set(0)

# button frame
google_image_button_frame = tk.Frame(google_image_frame)
# checkbox frame
download_google_image_checkbutton = tk.Frame(google_image_frame)
# download images checkbox
download_google_images = tk.Checkbutton(download_google_image_checkbutton, text="Download Images", font=MAIN_FONT, variable=download_google_imgs, onvalue=1, offvalue=0, command=lambda: gen_local.set(0))
# google search button
google_image_search_button = tk.Button(google_image_button_frame, text="Google Search", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGoogleImageThread)
# back button
google_image_back_button = tk.Button(google_image_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
download_google_image_checkbutton.pack(side="left")
download_google_images.pack(side="left")
google_image_search_button.pack(side="left", padx=35)
google_image_back_button.pack(side="left", padx=30)
# place google image frame on main frame
google_image_button_frame.place(x=175, y=495)
download_google_image_checkbutton.place(x=170, y=150)

# place button frame on word analysis frame
google_image_button_frame.pack(side="bottom", pady=10)

#--------------------------------Wiki Link Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
wiki_link_frame = tk.Frame(root, width=1000, height=600)
wiki_link_frame.pack_propagate(False)
wiki_link_frame.grid(row=0, column=0, sticky='news')
 
#--------------------------------Wiki Link Buttons--------------------------------------------------------------------------------------
# button frame
wiki_link_gen_button_frame = tk.Frame(wiki_link_frame)

wiki_link_gen_button = tk.Button(wiki_link_gen_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGenerateWikiLinkThread)
wiki_link_back_button = tk.Button(wiki_link_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
wiki_link_gen_button.pack(side="left", padx=35)
wiki_link_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
wiki_link_gen_button_frame.pack(side="bottom", pady=10)

#--------------------------------Who Are My Pairs Frame-----------------------------------------------------------------------------------------------
#------------------------------------------------------------------------------------------------------------------------------------------------
who_are_my_pairs_frame = tk.Frame(root, width=1000, height=600)
who_are_my_pairs_frame.pack_propagate(False)
who_are_my_pairs_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Who Are My Pairs Buttons-----------------------------------------------------------------------------------
#button frame
who_are_my_pairs_gen_button_frame = tk.Frame(who_are_my_pairs_frame)

who_are_my_pairs_gen_button = tk.Button(who_are_my_pairs_gen_button_frame, text="Generate Pairs", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startGeneratePairsThread)
who_are_my_pairs_back_button = tk.Button(who_are_my_pairs_gen_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack button into button frame
who_are_my_pairs_gen_button.pack(side="left", padx=35)
who_are_my_pairs_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame?
who_are_my_pairs_gen_button_frame.pack(side="bottom", pady=10)

#--------------------------------Translation Package Frame-----------------------------------------------------------------------------------------------
#--------------------------------------------------------------------------------------------------------------------------------------------------
translation_package_frame = tk.Frame(root, width=1000, height=600)
translation_package_frame.pack_propagate(False)
translation_package_frame.grid(row=0, column=0, sticky='news')

#--------------------------------Translation Package Buttons-----------------------------------------------------------------------------------------------
# generate button frame
translation_package_button_frame = tk.Frame(translation_package_frame)
# generate button
translation_package_generate_button = tk.Button(translation_package_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startTranslationPackageThread)
# help button
translation_package_help_button = tk.Button(translation_package_button_frame, text="Help", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=launchHelpSite)
# upload button
translation_package_back_button = tk.Button(translation_package_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))
# pack buttons into button frame
translation_package_generate_button.pack(side="left", padx=35)
translation_package_help_button.pack(side="left")
translation_package_back_button.pack(side="left", padx=30)

# place button frame on word analysis frame
translation_package_button_frame.pack(side="bottom", pady=10)

#-------------------------------First Person Frame-------------------------------------------------------------------------------------------
first_person_frame = tk.Frame(root, width=1000, height=600)
first_person_frame.pack_propagate(False)
first_person_frame.grid(row=0, column=0, sticky='news')

#------------------------------------First Person Buttons-------------------------------------------------------------------------------------
first_person_button_frame = tk.Frame(first_person_frame)

first_person_generate_button = tk.Button(first_person_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startFirstPersonThread)
first_person_back_button = tk.Button(first_person_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
first_person_generate_button.pack(side="left", padx=35)
first_person_back_button.pack(side="left", padx=30)

# place button frame on <something>
first_person_button_frame.pack(side="bottom", pady=10)
#-------------------------------Word Search Frame-------------------------------------------------------------------------------------------
word_puzzle_frame = tk.Frame(root, width=1000, height=600)
word_puzzle_frame.pack_propagate(False)
word_puzzle_frame.grid(row=0, column=0, sticky='news')

#------------------------------------Word Search Buttons-------------------------------------------------------------------------------------
word_puzzle_button_frame = tk.Frame(word_puzzle_frame)

word_puzzle_generate_button = tk.Button(word_puzzle_button_frame, text="Generate", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=startWordPuzzleThread)
word_puzzle_back_button = tk.Button(word_puzzle_button_frame, text="Back", font=LABEL_FONT, width=25, height=1, bg="#007FFF", fg="#ffffff", command=lambda: raiseFrame('main_frame'))

# pack buttons into button frame
word_puzzle_generate_button.pack(side="left", padx=35)
word_puzzle_back_button.pack(side="left", padx=30)

# place button frame on <something>
word_puzzle_button_frame.pack(side="bottom", pady=10)

#--------------------------------Layout Radio Buttons-------------------------------------------------------------------------------------
# layout variable
layout = tk.IntVar()
layout.set(4)

# layout frame
layout_frame = tk.Frame(word_puzzle_frame)
# layout radio buttons
layout_radio4 = tk.Radiobutton(layout_frame, text="Puzzle on Left - Text on right - Single Page - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=4)
layout_radio1 = tk.Radiobutton(layout_frame, text="Puzzle on Left Page - Text on Right Page - Two Page Mode - Portrait Mode",
                                font=MAIN_FONT, variable=layout, value=1)
layout_radio2 = tk.Radiobutton(layout_frame, text="Puzzle on Right - Text on Left - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=2)
layout_radio3 = tk.Radiobutton(layout_frame, text="Puzzle on Left - Text on Right - Single Page - Landscape Mode",
                                font=MAIN_FONT, variable=layout, value=3)
# pack radio buttons into layout frame
layout_radio4.pack(anchor="nw")
layout_radio1.pack(anchor="nw")
layout_radio2.pack(anchor="nw")
layout_radio3.pack(anchor="nw")
# place layout frame on main frame
layout_frame.place(x=175, y=150, width=800)
# layout radio buttons label
layout_label = tk.Label(word_puzzle_frame, text="Layout:", font=LABEL_FONT)
layout_label.place(x=25, y=170)

# separator line
separator1 = ttk.Separator(word_puzzle_frame)
separator1.place(x=175, y=150, relwidth=.8)

#--------------------------------Sort Radio Buttons---------------------------------------------------------------------------------------
# sort variable
sort_order = tk.IntVar()
sort_order.set(1)

# sort frame
sort_frame = tk.Frame(word_puzzle_frame)
# sort radio buttons
sort_radio1 = tk.Radiobutton(sort_frame, text="By Name", font=MAIN_FONT, variable=sort_order, value=1)
sort_radio2 = tk.Radiobutton(sort_frame, text="By ID", font=MAIN_FONT, variable=sort_order, value=2)
sort_radio3 = tk.Radiobutton(sort_frame, text="By Input Order", font=MAIN_FONT, variable=sort_order, value=3)
# pack radio buttons into sort frame
sort_radio1.pack(side="left")
sort_radio2.pack(side="left")
sort_radio3.pack(side="left")
# place sort frame on main frame
sort_frame.place(x=175, y=265, width=800)

# sort radio buttons label
sort_label = tk.Label(word_puzzle_frame, text="Sort Order:", font=LABEL_FONT)
sort_label.place(x=25, y=265)


# separator line
separator2 = ttk.Separator(word_puzzle_frame)
separator2.place(x=175, y=265, relwidth=.8)

#--------------------------------Word Radio Buttons---------------------------------------------------------------------------------------

#--------------------------------Preferences----------------------------------------------------------------------------------------------
preferences = {
    "WORD_COUNT": "10",
    "PUZ_WIDTH": "20",
}

# Using StringVar for dynamic updates
word_count_var = tk.StringVar(value=preferences["WORD_COUNT"])
puz_width_var = tk.StringVar(value=preferences["PUZ_WIDTH"])

# Frame for preference settings
preferences_frame = tk.Frame(word_puzzle_frame)
preferences_frame.place(x=175, y=295, width=800)

# Labels and entries for preferences
word_count_label = tk.Label(preferences_frame, text="Word Count:", font=MAIN_FONT)
word_count_label.grid(row=1, column=1, pady=10)
word_count = tk.Entry(preferences_frame, width=6, textvariable=word_count_var, font=MAIN_FONT)
word_count.grid(row=1, column=2)

puz_width_label = tk.Label(preferences_frame, text="Puzzle Width:", font=MAIN_FONT)
puz_width_label.grid(row=2, column=1, pady=10)
puz_width = tk.Entry(preferences_frame, width=6, textvariable=puz_width_var, font=MAIN_FONT)
puz_width.grid(row=2, column=2)


# Function to apply changes to preferences
def apply_changes():
    preferences["WORD_COUNT"] = word_count_var.get()
    preferences["PUZ_WIDTH"] = puz_width_var.get()
    print("Updated preferences:", preferences)

# Apply button
apply_button = tk.Button(preferences_frame, text="Apply Changes", command=apply_changes, font=LABEL_FONT)
apply_button.grid(row=4, column=1, columnspan=2, pady=10)

# preferences label
preferences_label = tk.Label(word_puzzle_frame, text="Preferences:", font=LABEL_FONT)
preferences_label.place(x=25, y=350)

# separator line
separator3 = ttk.Separator(word_puzzle_frame)
separator3.place(x=175, y=295, relwidth=.8)


#-------------------------------Start Main Frame----------------------------------------------------------------------------------------------

# raise main_frame to start
main_frame.tkraise()


# main gui loop
root.mainloop()
