import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import mysql.connector
import tkinter as tk
from tkinter import font
from tkinter import ttk
import random
import os

# Create GUI
def button_click():
    global input_topic, input_count, show_answers
    input_topic = topic_field.get()
    input_count = count_field.get()
    show_answers = show_answers.get()
    label.config(text="Input retrieved.")
    window.destroy()


window = tk.Tk()
custom_font = font.Font(size = 15)
window.geometry("800x800")
window.title("Quiz Master")
window.configure(bg = "gray32")
label = tk.Label(window, text="Welcome to QuizMaster", font = ("Arial", 40), fg = "dark orange", bg = "gray32")
topic_label = tk.Label(window, text="Please enter a topic:", font = ("Arial", 25), fg = "green yellow", bg = "gray32")
count_label = tk.Label(window, text="How many questions should be chosen?", font = ("Arial", 25), fg = "green yellow", bg = "gray32")
topic_field = tk.Entry(window, font = custom_font, width = 15)
count_field = tk.Entry(window, font = custom_font, width = 15)

show_answers = tk.BooleanVar()
answers_label = tk.Label(window, text = "Would you like to generate an answers slide?", font = ("Arial", 25), fg = "green yellow", bg = "gray32")
style = ttk.Style()
style.configure("custom.TRadiobutton", font = ("Arial", 15), foreground = "turquoise", background = "gray32")
answers_yes = ttk.Radiobutton(window, text = "Yes", variable = show_answers, value = True, style = "custom.TRadiobutton")
answers_no = ttk.Radiobutton(window, text = "No", variable = show_answers, value = False, style = "custom.TRadiobutton")
button = tk.Button(window, text="Submit", command=button_click, font = ("Arial", 25), width = 8, fg = "green yellow", bg = "gray16")

label.pack(pady = 50)
topic_label.pack(pady = 5)
topic_field.pack(pady = 5)
count_label.pack(pady = 10)
count_field.pack(pady = 10)
answers_label.pack(pady = 10)
answers_yes.pack()
answers_no.pack()
button.pack(pady = 50)

window.mainloop()

# Connect to db and run query
conn = mysql.connector.connect(
    host='localhost',
    user='root',
    #password='password',
    database='quiz_master'
)

cursor = conn.cursor()
query = "SELECT * FROM questions WHERE topic = '" + input_topic + "'"
cursor.execute(query)

results = cursor.fetchall()
if(input_count == ''):
     input_count = len(results)
input_count = int(input_count)
if (len(results) < input_count):
    input_count = len(results)


random_results = random.sample(results, input_count)
answers_list = []

# Create a Presentation object
presentation = Presentation()
slide_layout = presentation.slide_layouts[1]
slide_count = 0

for row in random_results:
        slide = presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        answers_list.append(row[7])

# Set the question as the slide title
        title_shape = shapes.title
        question = row[2]
        title_shape.text = question

# Add the potential answers in a textbox
        answers = [row[3], row[4], row[5], row[6]]
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "A. " + answers[0]
        p = tf.add_paragraph()
        p.text = "B. " + answers[1]
        p = tf.add_paragraph()
        p.text = "C. " + answers[2]
        p = tf.add_paragraph()
        p.text = "D. " + answers[3]
        slide_count += 1

# Add the code snippet as an image
        img_path = "D:/MyXampp/htdocs/QuizMasterDB/" + row[8]
        if os.path.exists(img_path):
            try:
                left = Inches(1)
                top = Inches(5)
                width = Inches(2)
                height = Inches(2)
                pic = slide.shapes.add_picture(img_path, left, top, width, height)
            except Exception as e:
                 print("ID: " + str(row[0]) + " Error adding image to slide.")
        else:
             print("ID: " + str(row[0]) + " Image not found.")

# Generate answers slide if desired
if show_answers:
    slide = presentation.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Answer Key"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    index = 1
    for answer in answers_list:
        p = tf.add_paragraph()
        p.font.size = Pt(15)
        p.text = str(index) + ". " + answer
        if (index % 10 == 0 and len(answers_list) > index):
            slide = presentation.slides.add_slide(slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = "Answer Key"
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
        index += 1


# Save the PowerPoint presentation to a file
presentation.save("/Users/Laith/Documents/images/fp4test.pptx")

cursor.close()
conn.close()